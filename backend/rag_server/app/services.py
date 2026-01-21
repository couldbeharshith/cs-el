"""
Unified Services Module for Insurance RAG API

This module consolidates all service classes into a single, well-organized file
"""

import os
import re
import logging
import asyncio
import hashlib
import itertools
import aiohttp
import tempfile
import time
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from urllib.parse import urlparse
from datetime import datetime

# Document processing imports
import fitz  # PyMuPDF
from docx import Document
from email import policy
from email.parser import BytesParser

# PPTX processing imports
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor
import pytesseract
from PIL import Image
import io

# XLSX processing imports
from openpyxl import load_workbook

# LLM and vector database imports
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold
from google import genai as google_genai
from google.genai import types

from groq import Groq
from pinecone import Pinecone

# Internal imports
from app.config import Config
from app.models import (
    DocumentChunk,
    ChunkMetadata,
    RelevantChunk,
    DocumentProcessingResult,
    VectorSearchResult
)
from app.pymupdf_processor import PyMuPDFProcessor, PyMuPDFError
from app.tiktoken_chunker import TikTokenChunker, TikTokenError, create_fallback_chunker


logger = logging.getLogger(__name__)


def chunks(iterable, batch_size=90):
    """A helper function to break an iterable into chunks of size batch_size."""
    it = iter(iterable)
    chunk = tuple(itertools.islice(it, batch_size))
    while chunk:
        yield chunk
        chunk = tuple(itertools.islice(it, batch_size))


class LLMKeyManager:
    """
    Simple cycling key manager that rotates through all keys without health tracking.
    Always cycles through all keys one by one, never marking any as failed.
    """

    def __init__(self, service_type: str, api_keys: List[str]):
        """
        Initialize key manager for a specific service type.

        Args:
            service_type: "gemini" or "groq"
            api_keys: List of API keys for the service
        """
        self.service_type = service_type.lower()
        self.api_keys = [key for key in api_keys if key and key.strip()]
        self.current_key_index = 0

        logger.info(
            f"LLMKeyManager initialized for {service_type} with {len(self.api_keys)} keys"
        )

    async def get_next_key(self) -> str:
        """
        Get the next API key in rotation.

        Returns:
            API key string

        Raises:
            Exception: If no API keys are configured
        """
        if not self.api_keys:
            raise Exception(f"No API keys configured for {self.service_type}")

        # Get next key in rotation
        key = self.api_keys[self.current_key_index]
        self.current_key_index = (self.current_key_index + 1) % len(self.api_keys)

        return key


class DocumentService:
    """
    Consolidated document processing service that handles PDF, DOCX, TXT, and EML files.
    Provides async document fetching from URLs and content extraction with chunking.
    """

    def __init__(self):
        self.supported_extensions = [
            ".pdf",
            ".docx",
            ".txt",
            ".eml",
            ".pptx",
            ".xlsx",
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".zip",
        ]
        # No file size limits - removed max_file_size restriction

        # Initialize PyMuPDF processor with markdown support
        self.pymupdf_processor = PyMuPDFProcessor(
            use_markdown=Config.USE_MARKDOWN_EXTRACTION
        )

        # Initialize TikToken chunker with fallback
        try:
            self.tiktoken_chunker = TikTokenChunker.create_from_config()
            self.chunking_fallback = None
            logger.info(
                f"TikTokenChunker initialized with encoding: {Config.TIKTOKEN_ENCODING}"
            )
        except TikTokenError as e:
            logger.warning(
                f"TikToken initialization failed: {e}, using fallback chunker"
            )
            self.tiktoken_chunker = None
            self.chunking_fallback = create_fallback_chunker(
                chunk_size=Config.TIKTOKEN_CHUNK_SIZE,
                chunk_overlap=Config.TIKTOKEN_CHUNK_OVERLAP,
            )

        logger.info(
            f"DocumentService initialized with PyMuPDF support for PDFs, DOCX, TXT, EML"
        )

    async def fetch_document(self, url: str) -> Tuple[bytes, str]:
        """
        Asynchronously fetch document from URL.

        Args:
            url: Document URL to fetch

        Returns:
            Tuple of (document_content_bytes, detected_content_type)

        Raises:
            Exception: If document fetching fails
        """
        try:
            logger.info(f"Fetching document from URL")

            async with aiohttp.ClientSession(
                timeout=aiohttp.ClientTimeout(total=None)
            ) as session:
                async with session.get(url) as response:
                    if response.status != 200:
                        raise Exception(
                            f"HTTP {response.status}: Failed to fetch document"
                        )

                    content_type = response.headers.get("content-type", "").lower()
                    content_length = response.headers.get("content-length")

                    # Check file size limit from config
                    max_size_bytes = (
                        Config.MAX_DOCUMENT_SIZE_MB * 1024 * 1024
                    )  # Convert MB to bytes

                    if content_length:
                        file_size_bytes = int(content_length)
                        file_size_mb = file_size_bytes / (1024 * 1024)

                        logger.info(
                            f"Document size: {file_size_mb:.2f} MB ({file_size_bytes} bytes)"
                        )

                        if file_size_bytes > max_size_bytes:
                            logger.warning(
                                f"Document size ({file_size_mb:.2f} MB) exceeds maximum allowed size ({Config.MAX_DOCUMENT_SIZE_MB} MB)"
                            )
                            # Return special indicator for oversized files instead of raising exception
                            return b"__FILE_TOO_LARGE__", "application/oversized"

                    content = await response.read()

                    # Double-check actual downloaded size
                    actual_size_bytes = len(content)
                    actual_size_mb = actual_size_bytes / (1024 * 1024)

                    if actual_size_bytes > max_size_bytes:
                        logger.warning(
                            f"Downloaded document size ({actual_size_mb:.2f} MB) exceeds maximum allowed size ({Config.MAX_DOCUMENT_SIZE_MB} MB)"
                        )
                        # Return special indicator for oversized files instead of raising exception
                        return b"__FILE_TOO_LARGE__", "application/oversized"

                    logger.info(
                        f"Successfully fetched document: {actual_size_mb:.2f} MB ({actual_size_bytes}) bytes, type: {content_type}"
                    )
                    return content, content_type

        except Exception as e:
            logger.error(f"Failed to fetch document from {url}: {e}")
            raise

    def detect_document_type(
        self, content: bytes, content_type: str = None, filename: str = None
    ) -> str:
        """
        Detect document type from content, content-type header, or filename.

        Args:
            content: Document content bytes
            content_type: HTTP content-type header
            filename: Original filename if available

        Returns:
            Document type: 'pdf', 'docx', 'pptx', 'xlsx', 'image', 'zip', 'txt', or 'eml'
        """
        # Check content-type header first
        if content_type:
            if "text" in content_type or "html" in content_type:
                return "txt"
            elif "pdf" in content_type:
                return "pdf"
            elif "wordprocessingml" in content_type or "msword" in content_type:
                return "docx"
            elif "presentationml" in content_type or "powerpoint" in content_type:
                return "pptx"
            elif "spreadsheetml" in content_type or "excel" in content_type:
                return "xlsx"
            elif content_type.startswith("image/"):
                return "image"
            elif "zip" in content_type or "application/zip" in content_type:
                return "zip"
            elif "text/plain" in content_type:
                return "txt"
            elif "message/rfc822" in content_type:
                return "eml"

        # Check filename extension
        if filename:
            ext = Path(filename).suffix.lower()
            if ext == ".pdf":
                return "pdf"
            elif ext == ".docx":
                return "docx"
            elif ext == ".pptx":
                return "pptx"
            elif ext == ".xlsx":
                return "xlsx"
            elif ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"]:
                return "image"
            elif ext == ".zip":
                return "zip"
            elif ext in [".txt", ".text"]:
                return "txt"
            elif ext == ".eml":
                return "eml"

        # Check content magic bytes
        if content.startswith(b"%PDF"):
            return "pdf"
        elif content.startswith(b"PK\x03\x04"):
            # ZIP-based formats - need to check internal structure
            if b"word/" in content:
                return "docx"
            elif b"ppt/" in content:
                return "pptx"
            elif b"xl/" in content:
                return "xlsx"
            else:
                return "zip"
        elif content.startswith(
            (b"\xff\xd8\xff", b"\x89PNG\r\n\x1a\n", b"GIF87a", b"GIF89a", b"BM")
        ):
            # JPEG, PNG, GIF, BMP magic bytes
            return "image"
        elif b"From:" in content[:1000] or b"To:" in content[:1000]:
            return "eml"
        else:
            return "txt"

    def _is_gemini_supported_file_type(self, content_type: str, filename: str) -> bool:
        """
        Check if a file type is supported by Gemini's direct upload API.

        Args:
            content_type: MIME type of the file
            filename: Name of the file

        Returns:
            True if the file type is supported by Gemini direct upload, False otherwise
        """
        # Gemini-supported MIME types for direct upload
        supported_mime_types = {
            # Text formats
            "text/plain",
            "text/html",
            "text/css",
            "text/javascript",
            "text/xml",
            "application/json",
            "application/xml",
            # PDF
            "application/pdf",
            # Images
            "image/jpeg",
            "image/png",
            "image/gif",
            "image/webp",
            "image/heic",
            "image/heif",
            # Audio
            "audio/wav",
            "audio/mp3",
            "audio/aiff",
            "audio/aac",
            "audio/ogg",
            "audio/flac",
            # Video
            "video/mp4",
            "video/mpeg",
            "video/mov",
            "video/avi",
            "video/x-flv",
            "video/mpg",
            "video/webm",
            "video/wmv",
            "video/3gpp",
        }

        # Handle content types with charset or other parameters
        base_content_type = content_type.split(";")[0].strip() if content_type else ""

        # Check exact match first
        if (
            content_type in supported_mime_types
            or base_content_type in supported_mime_types
        ):
            return True

        # Handle text/* types generically (most text types are supported)
        if base_content_type.startswith("text/"):
            return True

        # Check by file extension as fallback
        extension = Path(filename).suffix.lower()
        supported_extensions = {
            # Text formats
            ".txt",
            ".html",
            ".htm",
            ".css",
            ".js",
            ".json",
            ".xml",
            ".csv",
            ".md",
            ".log",
            # PDF
            ".pdf",
            # Images
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".webp",
            ".heic",
            ".heif",
            # Audio
            ".wav",
            ".mp3",
            ".aiff",
            ".aac",
            ".ogg",
            ".flac",
            # Video
            ".mp4",
            ".mpeg",
            ".mov",
            ".avi",
            ".flv",
            ".mpg",
            ".webm",
            ".wmv",
            ".3gp",
        }

        if extension in supported_extensions:
            return True

        # Special case: if it looks like a web URL, assume it's HTML/text
        if not extension and (
            filename.startswith("http") or "get-secret-token" in filename
        ):
            return True

        # Default to False for unsupported types (like .docx, .pptx, .xlsx, etc.)
        logger.debug(
            f"File type not supported by Gemini direct upload: {content_type} ({filename})"
        )
        return False

    async def process_document(
        self, url: str, vector_service: "VectorService" = None
    ) -> DocumentProcessingResult:
        """
        Process document from URL with hash-based caching support.

        Args:
            url: Document URL to process
            vector_service: Optional VectorService instance for cache checking

        Returns:
            DocumentProcessingResult with chunks and metadata
        """
        start_time = datetime.utcnow()

        try:
            # Fetch document
            content, content_type = await self.fetch_document(url)

            # Check if file is too large
            if (
                content == b"__FILE_TOO_LARGE__"
                and content_type == "application/oversized"
            ):
                logger.info(f"File too large, returning special processing result")
                processing_time = (datetime.utcnow() - start_time).total_seconds()
                return DocumentProcessingResult(
                    document_hash="file_too_large",
                    namespace="file_too_large",
                    chunks=[],
                    processing_time=processing_time,
                    cached=False,
                )

            # First, detect document type to prioritize image processing
            filename = Path(urlparse(url).path).name or "document"
            doc_type = self.detect_document_type(content, content_type, filename)

            # Check for image files FIRST - always use image processing for images regardless of size
            if doc_type == "image":
                logger.info(
                    f"Image document detected: {filename} - using image processing (bypassing direct upload)"
                )

                # Generate document hash for identification
                doc_hash = self.generate_document_hash(content)
                processing_time = (datetime.utcnow() - start_time).total_seconds()

                # Create a placeholder chunk to indicate this is an image
                from app.models import DocumentChunk, ChunkMetadata

                placeholder_metadata = ChunkMetadata(
                    document_id=doc_hash,
                    page_number=1,
                    line_number=1,
                    chunk_number=1,
                    document_type="image",
                    file_name=filename,
                )
                placeholder_chunk = DocumentChunk(
                    content=f"Image document: {filename}", metadata=placeholder_metadata
                )

                return DocumentProcessingResult(
                    document_hash=doc_hash,
                    namespace=f"image_{doc_hash}",  # Special namespace for images
                    chunks=[placeholder_chunk],
                    processing_time=processing_time,
                    cached=False,
                )

            # Check file size and type for direct Gemini upload (non-images only)
            file_size_mb = len(content) / (1024 * 1024)  # Convert bytes to MB

            # Check if file type is supported by Gemini direct upload
            is_gemini_supported_type = self._is_gemini_supported_file_type(
                content_type, filename
            )

            if (
                file_size_mb < Config.DIRECT_GEMINI_UPLOAD_THRESHOLD_MB
                and is_gemini_supported_type
            ):
                logger.info(
                    f"File size ({file_size_mb:.2f}MB) is below threshold ({Config.DIRECT_GEMINI_UPLOAD_THRESHOLD_MB}MB) and type ({content_type}) is Gemini-supported, using direct Gemini upload for: {filename}"
                )

                # Generate document hash for identification
                doc_hash = self.generate_document_hash(content)
                namespace = f"direct_gemini_{doc_hash}"

                processing_time = (datetime.utcnow() - start_time).total_seconds()

                # Create a special result indicating direct Gemini processing
                from app.models import DocumentChunk, ChunkMetadata

                direct_metadata = ChunkMetadata(
                    document_id=doc_hash,
                    page_number=1,
                    line_number=1,
                    chunk_number=1,
                    document_type="direct_gemini",
                    file_name=filename,
                )
                direct_chunk = DocumentChunk(
                    content=f"Direct Gemini document: {filename} ({file_size_mb:.2f}MB)",
                    metadata=direct_metadata,
                )

                return DocumentProcessingResult(
                    document_hash=doc_hash,
                    namespace=namespace,
                    chunks=[direct_chunk],
                    processing_time=processing_time,
                    cached=False,
                    direct_gemini_content=content,  # Store the raw content for Gemini upload
                    direct_gemini_content_type=content_type,
                )
            else:
                # Log why direct upload was not used
                if file_size_mb >= Config.DIRECT_GEMINI_UPLOAD_THRESHOLD_MB:
                    logger.info(
                        f"File size ({file_size_mb:.2f}MB) exceeds threshold ({Config.DIRECT_GEMINI_UPLOAD_THRESHOLD_MB}MB), using traditional processing for: {filename}"
                    )
                if not is_gemini_supported_type:
                    logger.info(
                        f"File type ({content_type}) not supported by Gemini direct upload, using traditional processing for: {filename}"
                    )

            # Generate document hash for caching (SHA-256 with first 16 chars)
            doc_hash = self.generate_document_hash(content)
            namespace = self.generate_namespace_from_hash(doc_hash)

            logger.info(f"Document hash: {doc_hash}, namespace: {namespace}")

            # Check if document is already cached (if vector_service provided)
            cached = False
            if vector_service:
                cached = await vector_service.check_namespace_exists(namespace)
                if cached:
                    logger.info(f"Document already cached")
                    processing_time = (datetime.utcnow() - start_time).total_seconds()
                    return DocumentProcessingResult(
                        document_hash=doc_hash,
                        namespace=namespace,
                        chunks=[],  # Empty chunks since we're using cached version
                        processing_time=processing_time,
                        cached=True,
                    )

            # Check for ZIP files and reject them early
            if doc_type == "zip":
                logger.info(f"ZIP file detected, rejecting: {filename}")
                processing_time = (datetime.utcnow() - start_time).total_seconds()
                return DocumentProcessingResult(
                    document_hash="zip_rejected",
                    namespace="zip_rejected",
                    chunks=[],
                    processing_time=processing_time,
                    cached=False,
                )

            logger.info(f"Processing {doc_type} document: {filename}")

            # Save to temporary file for processing with proper async resource management
            temp_path = None
            try:
                # Create temporary file asynchronously
                with tempfile.NamedTemporaryFile(
                    suffix=f".{doc_type}", delete=False
                ) as temp_file:
                    temp_file.write(content)
                    temp_path = temp_file.name

                # Process based on document type with proper async error handling and timeout
                try:
                    # Process document without timeout restrictions for large files
                    chunks = await self._process_document_with_processor_selection(
                        temp_path, doc_type, doc_hash
                    )

                    # Apply chunking based on document type
                    if doc_type == "pptx":
                        # PPTX uses slide-based chunking, no additional chunking needed
                        final_chunks = chunks
                        logger.info(
                            f"PPTX slide-based chunking: {len(chunks)} chunks (no TikToken chunking applied)"
                        )
                    else:
                        # Apply TikToken-based chunking for other document types
                        final_chunks = await asyncio.to_thread(
                            self._apply_chunking, chunks
                        )

                    # Store chunks in vector database if vector_service is provided
                    if vector_service and final_chunks:
                        storage_success = await self.store_document_chunks(
                            final_chunks, namespace, vector_service
                        )
                        if not storage_success:
                            logger.warning(
                                f"Failed to store chunks in vector database for namespace {namespace}"
                            )

                    processing_time = (datetime.utcnow() - start_time).total_seconds()

                    # Log processing time for monitoring (no limits enforced)
                    if processing_time > 60:  # Log for monitoring purposes only
                        logger.info(
                            f"Large file processing time: {processing_time:.2f}s"
                        )

                    return DocumentProcessingResult(
                        document_hash=doc_hash,
                        namespace=namespace,
                        chunks=final_chunks,
                        processing_time=processing_time,
                        cached=False,
                    )

                # No timeout restrictions - removed asyncio.TimeoutError handling
                except asyncio.CancelledError:
                    logger.warning(f"Document processing cancelled for {url}")
                    raise
                except Exception as e:
                    logger.error(f"Document processing failed for {url}: {e}")
                    raise

            finally:
                # Clean up temporary file with proper async resource cleanup
                if temp_path:
                    try:
                        await asyncio.to_thread(os.unlink, temp_path)
                    except Exception as cleanup_error:
                        logger.warning(
                            f"Failed to cleanup temporary file {temp_path}: {cleanup_error}"
                        )

        except Exception as e:
            logger.error(f"Failed to process document from {url}: {e}")
            raise

    def generate_document_hash(self, content: bytes) -> str:
        """
        Generate a fast hash of document content using SHA-256.

        Args:
            content: Document content as bytes

        Returns:
            First 16 characters of SHA-256 hash for efficient namespace naming
        """
        if not content:
            return "empty_doc"

        # Use SHA-256 for consistent, fast hashing
        hash_obj = hashlib.sha256(content)
        full_hash = hash_obj.hexdigest()

        # Return first 16 characters for namespace (keeps it short and unique)
        doc_hash = full_hash[:16]
        logger.debug(f"Generated document hash: {doc_hash} from {len(content)} bytes")

        return doc_hash

    def generate_namespace_from_hash(self, doc_hash: str) -> str:
        """
        Generate Pinecone namespace from document hash.

        Args:
            doc_hash: Document hash string

        Returns:
            Namespace string in format 'doc_{hash}'
        """
        namespace = f"doc_{doc_hash}"
        logger.debug(f"Generated namespace: {namespace}")
        return namespace

    async def _process_document_with_processor_selection(
        self, file_path: str, doc_type: str, document_id: str
    ) -> List[DocumentChunk]:
        """
        Process document using PyMuPDF or legacy processors based on configuration.
        Implements configurable processor selection with performance optimization.

        Args:
            file_path: Path to the document file
            doc_type: Document type ('pdf', 'docx', 'pptx', 'xlsx', 'image', 'txt', 'eml')
            document_id: Unique identifier for the document

        Returns:
            List of DocumentChunk objects
        """
        start_time = time.time()

        try:
            # Determine processor selection based on configuration
            should_use_pymupdf = self._should_use_pymupdf_for_document(doc_type)

            if should_use_pymupdf:
                logger.info(f"Using PyMuPDF processor for {doc_type} document")
                chunks = await self._process_document_with_pymupdf(
                    file_path, doc_type, document_id
                )
            else:
                logger.info(f"Using legacy processor for {doc_type} document")
                chunks = await self._process_document_legacy(
                    file_path, doc_type, document_id
                )

            processing_time = time.time() - start_time
            logger.debug(
                f"Document processing completed in {processing_time:.2f}s, produced {len(chunks)} chunks"
            )

            return chunks

        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(
                f"Document processing failed for {doc_type} after {processing_time:.2f}s: {e}"
            )
            raise

    def _should_use_pymupdf_for_document(self, doc_type: str) -> bool:
        """
        Determine if PyMuPDF should be used for a specific document type.
        Hardcoded processor selection:
        - PDF, TXT: PyMuPDF (if supported)
        - DOCX, EML: Legacy processors (PyMuPDF doesn't support these)
        - PPTX: python-pptx (with OCR)
        - XLSX: openpyxl
        - Images: custom processor

        Args:
            doc_type: Document type to check

        Returns:
            True if PyMuPDF should be used, False otherwise
        """
        # Use PyMuPDF for these document types if it supports them
        if doc_type in ["pdf", "txt"] and self.pymupdf_processor.can_process(doc_type):
            return True

        # Use legacy/custom processors for all other document types
        return False

    async def _process_document_with_pymupdf(
        self, file_path: str, doc_type: str, document_id: str
    ) -> List[DocumentChunk]:
        """
        Process document using PyMuPDF with performance optimization and fallback to legacy processors.

        Args:
            file_path: Path to the document file
            doc_type: Document type
            document_id: Unique identifier for the document

        Returns:
            List of DocumentChunk objects
        """
        start_time = time.time()

        try:
            # Use PyMuPDF processor with performance monitoring
            chunks = await asyncio.to_thread(
                self.pymupdf_processor.process_document_to_chunks,
                file_path,
                document_id,
                doc_type,
            )

            processing_time = time.time() - start_time
            logger.info(
                f"PyMuPDF successfully processed {doc_type} document in {processing_time:.2f}s: {len(chunks)} chunks"
            )

            # Validate chunks
            if not chunks:
                logger.warning(
                    f"PyMuPDF produced no chunks for {doc_type}, falling back to legacy processor"
                )
                return await self._process_document_legacy(
                    file_path, doc_type, document_id
                )

            return chunks

        except PyMuPDFError as e:
            processing_time = time.time() - start_time
            logger.warning(
                f"PyMuPDF processing failed for {doc_type} after {processing_time:.2f}s: {e.message}, falling back to legacy processor"
            )
            # Fallback to legacy processor
            return await self._process_document_legacy(file_path, doc_type, document_id)

        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(
                f"PyMuPDF processing failed unexpectedly for {doc_type} after {processing_time:.2f}s: {e}, falling back to legacy processor"
            )
            # Fallback to legacy processor
            return await self._process_document_legacy(file_path, doc_type, document_id)

    async def _process_document_legacy(
        self, file_path: str, doc_type: str, document_id: str
    ) -> List[DocumentChunk]:
        """
        Process document using legacy processors with performance monitoring.

        Args:
            file_path: Path to the document file
            doc_type: Document type
            document_id: Unique identifier for the document

        Returns:
            List of DocumentChunk objects
        """
        start_time = time.time()

        try:
            if doc_type == "pdf":
                chunks = await asyncio.to_thread(
                    self._process_pdf, file_path, document_id
                )
            elif doc_type == "docx":
                chunks = await asyncio.to_thread(
                    self._process_docx, file_path, document_id
                )
            elif doc_type == "pptx":
                chunks = await self._process_pptx(file_path, document_id)
            elif doc_type == "xlsx":
                chunks = await asyncio.to_thread(
                    self._process_xlsx, file_path, document_id
                )
            elif doc_type == "image":
                chunks = await self._process_image(file_path, document_id)
            elif doc_type == "eml":
                chunks = await asyncio.to_thread(
                    self._process_eml, file_path, document_id
                )
            else:  # txt
                chunks = await asyncio.to_thread(
                    self._process_txt, file_path, document_id
                )

            logger.info(
                f"Legacy processor completed {doc_type} processing in {(time.time() - start_time):.2f}s: {len(chunks)} chunks"
            )

            return chunks

        except Exception as e:
            logger.error(
                f"Legacy processing failed for {doc_type} after {(time.time() - start_time):.2f}s: {e}"
            )
            raise

    def _process_pdf(self, file_path: str, document_id: str) -> List[DocumentChunk]:
        """Process PDF document and extract text chunks using PyMuPDF (fitz)."""
        chunks = []
        chunk_number = 1

        try:
            # Use PyMuPDF (fitz) for PDF processing
            pdf_document = fitz.open(file_path)

            for page_num in range(len(pdf_document)):
                try:
                    page = pdf_document.load_page(page_num)
                    page_text = page.get_text()

                    if page_text and page_text.strip():
                        # Clean and validate text
                        clean_text = self._clean_extracted_text(page_text)
                        if clean_text:
                            metadata = ChunkMetadata(
                                document_id=document_id,
                                page_number=page_num + 1,  # 1-based page numbering
                                line_number=1,
                                chunk_number=chunk_number,
                                document_type="pdf",
                                file_name=Path(file_path).name,
                            )

                            chunk = DocumentChunk(content=clean_text, metadata=metadata)
                            chunks.append(chunk)
                            chunk_number += 1
                except Exception as e:
                    logger.warning(
                        f"Failed to extract text from PDF page {page_num + 1}: {e}"
                    )
                    continue

            pdf_document.close()

        except Exception as e:
            logger.error(f"PyMuPDF (fitz) failed to process PDF: {e}")
            raise Exception(f"Failed to process PDF with PyMuPDF: {e}")

        if not chunks:
            raise Exception("No readable text found in PDF document")

        return chunks

    def _process_docx(self, file_path: str, document_id: str) -> List[DocumentChunk]:
        """Process DOCX document and extract text chunks."""
        chunks = []
        chunk_number = 1

        try:
            doc = Document(file_path)

            # Process paragraphs
            for para_idx, paragraph in enumerate(doc.paragraphs, 1):
                text = paragraph.text.strip()
                if text:
                    metadata = ChunkMetadata(
                        document_id=document_id,
                        page_number=para_idx,
                        line_number=1,
                        chunk_number=chunk_number,
                        document_type="docx",
                        file_name=Path(file_path).name,
                    )

                    chunk = DocumentChunk(content=text, metadata=metadata)
                    chunks.append(chunk)
                    chunk_number += 1

            # Process tables
            for table_idx, table in enumerate(doc.tables, 1):
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        table_text.append(" | ".join(row_text))

                if table_text:
                    combined_table_text = "\n".join(table_text)
                    metadata = ChunkMetadata(
                        document_id=document_id,
                        page_number=len(doc.paragraphs) + table_idx,
                        line_number=1,
                        chunk_number=chunk_number,
                        document_type="docx",
                        file_name=Path(file_path).name,
                    )

                    chunk = DocumentChunk(
                        content=combined_table_text, metadata=metadata
                    )
                    chunks.append(chunk)
                    chunk_number += 1

        except Exception as e:
            logger.error(f"Failed to process DOCX: {e}")
            raise Exception(f"Failed to process DOCX document: {e}")

        if not chunks:
            raise Exception("No readable content found in DOCX document")

        return chunks

    def _process_eml(self, file_path: str, document_id: str) -> List[DocumentChunk]:
        """Process EML email document and extract text chunks."""
        chunks = []
        chunk_number = 1

        try:
            with open(file_path, "rb") as file:
                parser = BytesParser(policy=policy.default)
                msg = parser.parse(file)

            # Extract headers
            headers = {
                "from": msg.get("From", ""),
                "to": msg.get("To", ""),
                "subject": msg.get("Subject", ""),
                "date": msg.get("Date", ""),
            }

            # Create header chunk
            header_content = []
            for key, value in headers.items():
                if value and value.strip():
                    header_content.append(f"{key.title()}: {value}")

            if header_content:
                header_text = "\n".join(header_content)
                metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=1,
                    line_number=1,
                    chunk_number=chunk_number,
                    document_type="eml",
                    file_name=Path(file_path).name,
                )

                chunk = DocumentChunk(content=header_text, metadata=metadata)
                chunks.append(chunk)
                chunk_number += 1

            # Extract body content
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_maintype() == "multipart":
                        continue

                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            text = payload.decode("utf-8", errors="ignore").strip()
                            if text:
                                metadata = ChunkMetadata(
                                    document_id=document_id,
                                    page_number=chunk_number,
                                    line_number=1,
                                    chunk_number=chunk_number,
                                    document_type="eml",
                                    file_name=Path(file_path).name,
                                )

                                chunk = DocumentChunk(content=text, metadata=metadata)
                                chunks.append(chunk)
                                chunk_number += 1
            else:
                # Single part message
                if msg.get_content_type() == "text/plain":
                    payload = msg.get_payload(decode=True)
                    if payload:
                        text = payload.decode("utf-8", errors="ignore").strip()
                        if text:
                            metadata = ChunkMetadata(
                                document_id=document_id,
                                page_number=chunk_number,
                                line_number=1,
                                chunk_number=chunk_number,
                                document_type="eml",
                                file_name=Path(file_path).name,
                            )

                            chunk = DocumentChunk(content=text, metadata=metadata)
                            chunks.append(chunk)
                            chunk_number += 1

        except Exception as e:
            logger.error(f"Failed to process EML: {e}")
            raise Exception(f"Failed to process EML document: {e}")

        if not chunks:
            raise Exception("No readable content found in EML document")

        return chunks

    def _process_txt(self, file_path: str, document_id: str) -> List[DocumentChunk]:
        """Process TXT document and extract text chunks."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
                content = file.read().strip()

            if not content:
                raise Exception("No readable content found in TXT document")

            metadata = ChunkMetadata(
                document_id=document_id,
                page_number=1,
                line_number=1,
                chunk_number=1,
                document_type="txt",
                file_name=Path(file_path).name,
            )

            chunk = DocumentChunk(content=content, metadata=metadata)
            return [chunk]

        except Exception as e:
            logger.error(f"Failed to process TXT: {e}")
            raise Exception(f"Failed to process TXT document: {e}")

    async def _process_pptx(
        self, file_path: str, document_id: str
    ) -> List[DocumentChunk]:
        """
        Process PPTX document with slide-based chunking using multithreading.
        Each slide becomes one chunk (unless > 1200 chars, then split).

        Args:
            file_path: Path to the PPTX file
            document_id: Unique identifier for the document

        Returns:
            List of DocumentChunk objects with slide-based content
        """
        chunks = []

        try:
            # Load presentation
            presentation = Presentation(file_path)
            total_slides = len(presentation.slides)

            if total_slides == 0:
                raise Exception("No slides found in PPTX document")

            logger.info(
                f"Processing PPTX with {total_slides} slides using multithreading (slide-based chunking)"
            )

            # Process slides in parallel using ThreadPoolExecutor
            from concurrent.futures import as_completed

            with ThreadPoolExecutor(
                max_workers=None
            ) as executor:  # Use default max_workers like the script
                futures = {
                    executor.submit(
                        self._process_single_slide_combined,
                        slide_idx + 1,
                        slide,
                        document_id,
                        file_path,
                    ): slide_idx
                    + 1
                    for slide_idx, slide in enumerate(presentation.slides)
                }

                # Collect results using as_completed like the working script
                # We need to maintain chunk numbering across all slides
                global_chunk_number = 1
                slide_results = {}

                for future in as_completed(futures):
                    slide_number = futures[future]
                    try:
                        slide_content = future.result()
                        if slide_content:
                            slide_results[slide_number] = slide_content
                    except Exception as e:
                        logger.error(f"Failed to process slide {slide_number}: {e}")
                        continue

                # Process slides in order to maintain sequential chunk numbering
                for slide_number in sorted(slide_results.keys()):
                    slide_content = slide_results[slide_number]
                    slide_chunks = self._create_slide_based_chunks(
                        slide_content,
                        slide_number,
                        document_id,
                        file_path,
                        global_chunk_number,
                    )
                    chunks.extend(slide_chunks)
                    global_chunk_number += len(slide_chunks)

            if not chunks:
                raise Exception("No readable content found in PPTX document")

            logger.info(
                f"Successfully processed PPTX: {len(chunks)} slide-based chunks from {total_slides} slides"
            )
            return chunks

        except Exception as e:
            logger.error(f"Failed to process PPTX: {e}")
            raise Exception(f"Failed to process PPTX document: {e}")

    def _process_single_slide_combined(
        self, slide_number: int, slide, document_id: str, file_path: str
    ) -> str:
        """
        Process a single slide and combine all content into one string.
        combines title, subtitle, text, tables, OCR.

        Args:
            slide_number: 1-based slide number
            slide: python-pptx Slide object
            document_id: Document identifier
            file_path: Path to the PPTX file

        Returns:
            Combined content string for the slide
        """
        import threading

        # Use thread-local storage for OCR to avoid conflicts
        if not hasattr(threading.current_thread(), "ocr_initialized"):
            threading.current_thread().ocr_initialized = True

        lines = [f"--- Slide {slide_number} ---"]

        try:
            # Title
            title_text = self._extract_slide_title(slide)
            if title_text:
                lines.append(f"Title: {title_text}")

            # Subtitle
            subtitle_text = self._extract_slide_subtitle(slide)
            if subtitle_text:
                lines.append(f"Subtitle: {subtitle_text}")

            # Text content (bullet points and other text)
            text_content = self._extract_slide_text_content(slide)
            if text_content:
                lines.append("Content:")
                lines.append(text_content)

            # Tables
            table_content = self._extract_slide_tables(slide)
            if table_content:
                lines.append("Tables:")
                lines.append(table_content)

            # Image OCR text
            image_ocr_content = self._extract_slide_image_ocr(slide)
            if image_ocr_content:
                lines.append("Image Text (OCR):")
                lines.append(image_ocr_content)

            # Combine all content
            combined_content = "\n".join(lines)

            return combined_content

        except Exception as e:
            logger.error(f"Failed to process slide {slide_number}: {e}")
            return f"--- Slide {slide_number} ---\nError processing slide: {str(e)}"

    def _create_slide_based_chunks(
        self,
        slide_content: str,
        slide_number: int,
        document_id: str,
        file_path: str,
        start_chunk_number: int = 1,
    ) -> List[DocumentChunk]:
        """
        Create chunks from slide content. One slide = one chunk unless > 1200 chars.

        Args:
            slide_content: Combined content from the slide
            slide_number: 1-based slide number
            document_id: Document identifier
            file_path: Path to the PPTX file
            start_chunk_number: Starting chunk number for sequential numbering

        Returns:
            List of DocumentChunk objects (1 chunk per slide, or multiple if slide > 1200 chars)
        """
        chunks = []

        if len(slide_content) <= 1200:
            # Slide fits in one chunk
            metadata = ChunkMetadata(
                document_id=document_id,
                page_number=slide_number,
                line_number=1,
                chunk_number=start_chunk_number,
                document_type="pptx",
                file_name=Path(file_path).name,
                slide_number=slide_number,
                content_type="slide",
            )
            chunks.append(DocumentChunk(content=slide_content, metadata=metadata))

        else:
            # Slide is too long, split into multiple chunks

            # Split content into chunks of ~1200 characters
            chunk_size = 1200
            overlap = 100  # Small overlap to maintain context

            start = 0
            chunk_num = start_chunk_number

            while start < len(slide_content):
                end = start + chunk_size

                # If this isn't the last chunk, try to break at a word boundary
                if end < len(slide_content):
                    # Look for a good break point (newline, period, or space)
                    for break_char in ["\n", ". ", " "]:
                        break_pos = slide_content.rfind(break_char, start, end)
                        if break_pos > start:
                            end = break_pos + len(break_char)
                            break

                chunk_text = slide_content[start:end].strip()

                if chunk_text:
                    metadata = ChunkMetadata(
                        document_id=document_id,
                        page_number=slide_number,
                        line_number=1,
                        chunk_number=chunk_num,
                        document_type="pptx",
                        file_name=Path(file_path).name,
                        slide_number=slide_number,
                        content_type="slide_part",
                    )
                    chunks.append(DocumentChunk(content=chunk_text, metadata=metadata))
                    chunk_num += 1

                # Move start position with overlap
                start = max(start + chunk_size - overlap, end)

            logger.info(f"Slide {slide_number}: Split into {len(chunks)} chunks")

        return chunks

    def _process_slide(
        self, slide, slide_number: int, document_id: str, file_path: str
    ) -> List[DocumentChunk]:
        """
        Process a single slide and extract all content types.

        Args:
            slide: python-pptx Slide object
            slide_number: 1-based slide number
            document_id: Document identifier
            file_path: Path to the PPTX file

        Returns:
            List of DocumentChunk objects for this slide
        """
        slide_chunks = []
        chunk_number = 1

        try:
            # Extract title
            title_text = self._extract_slide_title(slide)
            if title_text:
                metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=slide_number,
                    line_number=1,
                    chunk_number=chunk_number,
                    document_type="pptx",
                    file_name=Path(file_path).name,
                    slide_number=slide_number,
                    content_type="title",
                )
                slide_chunks.append(
                    DocumentChunk(content=title_text, metadata=metadata)
                )
                chunk_number += 1

            # Extract subtitle
            subtitle_text = self._extract_slide_subtitle(slide)
            if subtitle_text:
                metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=slide_number,
                    line_number=1,
                    chunk_number=chunk_number,
                    document_type="pptx",
                    file_name=Path(file_path).name,
                    slide_number=slide_number,
                    content_type="subtitle",
                )
                slide_chunks.append(
                    DocumentChunk(content=subtitle_text, metadata=metadata)
                )
                chunk_number += 1

            # Extract text content (bullet points and other text)
            text_content = self._extract_slide_text_content(slide)
            if text_content:
                metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=slide_number,
                    line_number=1,
                    chunk_number=chunk_number,
                    document_type="pptx",
                    file_name=Path(file_path).name,
                    slide_number=slide_number,
                    content_type="text",
                )
                slide_chunks.append(
                    DocumentChunk(content=text_content, metadata=metadata)
                )
                chunk_number += 1

            # Extract tables
            table_content = self._extract_slide_tables(slide)
            if table_content:
                metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=slide_number,
                    line_number=1,
                    chunk_number=chunk_number,
                    document_type="pptx",
                    file_name=Path(file_path).name,
                    slide_number=slide_number,
                    content_type="table",
                )
                slide_chunks.append(
                    DocumentChunk(content=table_content, metadata=metadata)
                )
                chunk_number += 1

            # Extract image OCR text with enhanced logging
            logger.info(f"Attempting OCR extraction for slide {slide_number}")
            image_ocr_content = self._extract_slide_image_ocr(slide)
            if image_ocr_content:
                logger.info(
                    f"OCR content extracted for slide {slide_number}: {len(image_ocr_content)} characters"
                )
                metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=slide_number,
                    line_number=1,
                    chunk_number=chunk_number,
                    document_type="pptx",
                    file_name=Path(file_path).name,
                    slide_number=slide_number,
                    content_type="image_ocr",
                )
                slide_chunks.append(
                    DocumentChunk(content=image_ocr_content, metadata=metadata)
                )
                chunk_number += 1
            else:
                logger.warning(f"No OCR content extracted for slide {slide_number}")

            return slide_chunks

        except Exception as e:
            logger.warning(f"Failed to process slide {slide_number}: {e}")
            return []

    def _extract_slide_title(self, slide) -> str:
        """Extract title from slide."""
        try:
            if hasattr(slide, "shapes") and slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                return title_text if title_text else ""
        except Exception as e:
            logger.debug(f"Failed to extract slide title: {e}")
        return ""

    def _extract_slide_subtitle(self, slide) -> str:
        """Extract subtitle from slide."""
        try:
            if hasattr(slide, "placeholders"):
                for placeholder in slide.placeholders:
                    if (
                        hasattr(placeholder, "placeholder_format")
                        and placeholder.placeholder_format
                    ):
                        if (
                            placeholder.placeholder_format.type == 2
                        ):  # Subtitle placeholder
                            subtitle_text = placeholder.text.strip()
                            return subtitle_text if subtitle_text else ""
        except Exception as e:
            logger.debug(f"Failed to extract slide subtitle: {e}")
        return ""

    def _extract_slide_text_content(self, slide) -> str:
        """Extract all text content from slide (excluding title and subtitle)."""
        text_content = []

        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Skip title and subtitle shapes
                    if hasattr(slide.shapes, "title") and shape == slide.shapes.title:
                        continue

                    # Check if it's a subtitle placeholder
                    is_subtitle = False
                    if (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format
                    ):
                        if shape.placeholder_format.type == 2:  # Subtitle
                            is_subtitle = True

                    if not is_subtitle:
                        # Handle text frames with paragraphs for proper bullet point formatting
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            formatted_text = self._format_text_frame(shape.text_frame)
                            if formatted_text:
                                text_content.append(formatted_text)
                        else:
                            clean_text = shape.text.strip()
                            if clean_text:
                                text_content.append(clean_text)

        except Exception as e:
            logger.debug(f"Failed to extract slide text content: {e}")

        return "\n\n".join(text_content) if text_content else ""

    def _format_text_frame(self, text_frame) -> str:
        """Format text frame with proper bullet point indentation."""
        formatted_lines = []

        try:
            for paragraph in text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Get indentation level (0-based)
                    level = paragraph.level if hasattr(paragraph, "level") else 0
                    indent = "  " * level  # 2 spaces per level

                    # Add bullet point for indented items
                    if level > 0:
                        formatted_lines.append(f"{indent}- {text}")
                    else:
                        formatted_lines.append(text)

        except Exception as e:
            logger.debug(f"Failed to format text frame: {e}")
            # Fallback to simple text extraction
            try:
                return text_frame.text.strip()
            except:
                return ""

        return "\n".join(formatted_lines)

    def _extract_slide_tables(self, slide) -> str:
        """Extract tables from slide and convert to markdown format."""
        table_content = []

        try:
            for shape in slide.shapes:
                if hasattr(shape, "table"):
                    table = shape.table
                    markdown_table = self._convert_table_to_markdown(table)
                    if markdown_table:
                        table_content.append(markdown_table)

        except Exception as e:
            logger.debug(f"Failed to extract slide tables: {e}")

        return "\n\n".join(table_content) if table_content else ""

    def _convert_table_to_markdown(self, table) -> str:
        """Convert PPTX table to markdown format."""
        try:
            rows = []

            # Extract table data
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    # Escape pipe characters in cell content
                    cell_text = cell_text.replace("|", "\\|")
                    row_data.append(cell_text)
                rows.append(row_data)

            if not rows:
                return ""

            # Create markdown table
            markdown_lines = []

            # Header row
            if rows:
                header = "| " + " | ".join(rows[0]) + " |"
                markdown_lines.append(header)

                # Separator row
                separator = "| " + " | ".join(["---"] * len(rows[0])) + " |"
                markdown_lines.append(separator)

                # Data rows
                for row in rows[1:]:
                    data_row = "| " + " | ".join(row) + " |"
                    markdown_lines.append(data_row)

            return "\n".join(markdown_lines)

        except Exception as e:
            logger.debug(f"Failed to convert table to markdown: {e}")
            return ""

    def _extract_slide_image_ocr(self, slide) -> str:
        """Extract text from images in slide using OCR with multithreading - based on working script."""
        ocr_content = []

        try:
            # Import MSO_SHAPE_TYPE for proper image detection
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from concurrent.futures import ThreadPoolExecutor, as_completed

            # Collect all images first
            images_to_process = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        images_to_process.append(img)
                    except Exception as e:
                        logger.warning(f"Failed to extract image from shape: {e}")
                        continue

            if not images_to_process:
                return ""

            # Process images in parallel using ThreadPoolExecutor
            with ThreadPoolExecutor(
                max_workers=min(4, len(images_to_process))
            ) as executor:
                # Submit all OCR tasks
                futures = {
                    executor.submit(self._perform_ocr_with_fallback, img): idx
                    for idx, img in enumerate(images_to_process)
                }

                # Collect results as they complete
                results = {}
                for future in as_completed(futures):
                    img_idx = futures[future]
                    try:
                        ocr_text = future.result()
                        if ocr_text:
                            results[img_idx] = ocr_text
                    except Exception as e:
                        logger.warning(f"OCR failed for image {img_idx}: {e}")
                        continue

                # Add results in order
                for idx in sorted(results.keys()):
                    ocr_content.append(results[idx])

        except Exception as e:
            logger.error(f"Failed to extract slide image OCR: {e}")

        return "\n\n".join(ocr_content) if ocr_content else ""

    def _perform_ocr_with_fallback(self, image: Image.Image) -> str:
        """Perform OCR on image - thread-safe implementation."""
        try:
            # Image is already converted to RGB in the calling function
            logger.info(f"Attempting OCR on image: {image.size}, mode: {image.mode}")

            ocr_txt = pytesseract.image_to_string(image).strip().replace("\n", " ")
            return ocr_txt if ocr_txt else ""

        except pytesseract.TesseractNotFoundError:
            logger.error("⚠️  Tesseract not found—skipping OCR.")
            return ""
        except ImportError as e:
            logger.error(f"pytesseract not available: {e}")
            return ""
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""

    def _clean_extracted_text(self, text: str) -> str:
        """Clean extracted text by removing PDF metadata and formatting artifacts."""
        if not text:
            return ""

        # Check for PDF metadata patterns
        pdf_metadata_patterns = [
            "%PDF",
            "endobj",
            "/Type/",
            "obj\n",
            "<<",
            ">>",
            "stream",
            "endstream",
        ]

        if any(pattern in text for pattern in pdf_metadata_patterns):
            logger.warning("Text contains PDF metadata, skipping")
            return ""

        # Clean up whitespace and formatting
        lines = text.split("\n")
        clean_lines = [line.strip() for line in lines if line.strip()]

        return "\n".join(clean_lines)

    def _apply_chunking(
        self,
        chunks: List[DocumentChunk],
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
    ) -> List[DocumentChunk]:
        """
        Apply TikToken-based text chunking to document chunks with fallback support.
        Optimized for large files with efficient processing.

        Args:
            chunks: List of DocumentChunk objects to process
            chunk_size: Optional custom chunk size (defaults to Config.TIKTOKEN_CHUNK_SIZE)
            chunk_overlap: Optional custom chunk overlap (defaults to Config.TIKTOKEN_CHUNK_OVERLAP)

        Returns:
            List of properly sized DocumentChunk objects with token-aware chunking
        """
        if not chunks:
            return []

        final_chunks = []

        # Use provided parameters or fall back to config defaults
        effective_chunk_size = chunk_size or Config.TIKTOKEN_CHUNK_SIZE
        effective_overlap = chunk_overlap or Config.TIKTOKEN_CHUNK_OVERLAP

        logger.info(
            f"Applying TikToken chunking to {len(chunks)} chunks with size={effective_chunk_size}, overlap={effective_overlap}"
        )

        # Process chunks efficiently for large files
        total_content_size = sum(len(chunk.content) for chunk in chunks)
        logger.info(f"Processing total content size: {total_content_size:,} characters")

        for i, chunk in enumerate(chunks):
            try:
                # Log progress for large files
                if i > 0 and i % 100 == 0:
                    logger.info(
                        f"Processed {i}/{len(chunks)} chunks ({(i/len(chunks)*100):.1f}%)"
                    )

                # Use TikToken chunker if available
                if self.tiktoken_chunker:
                    # Check if we need to chunk based on token count
                    token_count = self.tiktoken_chunker.count_tokens(chunk.content)

                    if token_count <= effective_chunk_size:
                        # Content fits in one chunk, enhance metadata with token info
                        enhanced_metadata = (
                            self.tiktoken_chunker._enhance_metadata_with_tokens(
                                chunk.metadata,
                                0,
                                token_count,
                                0,
                                len(chunk.content),
                                token_count,
                            )
                        )
                        enhanced_chunk = DocumentChunk(
                            content=chunk.content, metadata=enhanced_metadata
                        )
                        final_chunks.append(enhanced_chunk)
                    else:
                        # Content needs chunking
                        sub_chunks = self.tiktoken_chunker.chunk_text(
                            chunk.content, chunk.metadata
                        )
                        final_chunks.extend(sub_chunks)

                        logger.debug(
                            f"TikToken chunked large content ({token_count} tokens) into {len(sub_chunks)} chunks"
                        )

                # Fallback to character-based chunking if TikToken is not available
                elif self.chunking_fallback:
                    if len(chunk.content) <= effective_chunk_size:
                        final_chunks.append(chunk)
                    else:
                        sub_chunks = self.chunking_fallback(
                            chunk.content, chunk.metadata
                        )
                        final_chunks.extend(sub_chunks)

                        logger.debug(
                            f"Fallback chunked large content ({len(chunk.content)} chars) into {len(sub_chunks)} chunks"
                        )

                else:
                    # No chunking available, use original chunk
                    logger.warning("No chunking method available, using original chunk")
                    final_chunks.append(chunk)

            except TikTokenError as e:
                logger.warning(
                    f"TikToken chunking failed for chunk {i}: {e}, using fallback"
                )
                # Use fallback chunker
                if self.chunking_fallback:
                    if len(chunk.content) <= effective_chunk_size:
                        final_chunks.append(chunk)
                    else:
                        sub_chunks = self.chunking_fallback(
                            chunk.content, chunk.metadata
                        )
                        final_chunks.extend(sub_chunks)
                else:
                    final_chunks.append(chunk)

            except Exception as e:
                logger.error(
                    f"Unexpected error during chunking for chunk {i}: {e}, using original chunk"
                )
                final_chunks.append(chunk)

        logger.info(
            f"Chunking complete: {len(chunks)} original chunks -> {len(final_chunks)} final chunks"
        )
        return final_chunks

    async def process_document_from_url(
        self, url: str, vector_service: "VectorService" = None
    ) -> DocumentProcessingResult:
        """
        Process document from URL and return result with namespace information.

        Args:
            url: Document URL to process
            vector_service: Optional VectorService instance to reuse (avoids creating new instances)

        Returns:
            DocumentProcessingResult with document hash, namespace, and chunks
        """
        try:
            # Use provided vector service or create new one if not provided
            if vector_service is None:
                # Import VectorService here to avoid circular import
                from app.services import VectorService

                vector_service = VectorService()
                logger.warning(
                    "Creating new VectorService instance - consider passing existing instance"
                )

            # Use the existing process_document method
            result = await self.process_document(url, vector_service)
            return result

        except Exception as e:
            logger.error(f"Failed to process document from URL {url}: {e}")
            raise

    def prepare_chunks_for_vector_storage(
        self, chunks: List[DocumentChunk], document_hash: str
    ) -> List[Dict[str, Any]]:
        """
        Prepare document chunks for vector storage following pinetest.py patterns.

        Args:
            chunks: List of DocumentChunk objects
            document_hash: Document hash for unique ID generation

        Returns:
            List of records formatted for Pinecone upsert following pinetest.py format:
            {"_id": "unique-id", "text": "content", "metadata": "fields"}
        """
        vector_records = []

        for chunk in chunks:
            # Generate unique ID using document hash and chunk number
            unique_id = f"{document_hash}_{chunk.metadata.chunk_number}"

            # Prepare metadata fields for storage
            metadata_fields = {
                "document_id": chunk.metadata.document_id,
                "page_number": chunk.metadata.page_number,
                "line_number": chunk.metadata.line_number,
                "chunk_number": chunk.metadata.chunk_number,
                "document_type": chunk.metadata.document_type,
                "file_name": chunk.metadata.file_name,
                "document_hash": document_hash,
            }

            # Create record in pinetest.py format - metadata fields are stored directly
            record = {
                "_id": unique_id,
                "text": chunk.content,
                "document_id": metadata_fields["document_id"],
                "page_number": metadata_fields["page_number"],
                "line_number": metadata_fields["line_number"],
                "chunk_number": metadata_fields["chunk_number"],
                "document_type": metadata_fields["document_type"],
                "file_name": metadata_fields["file_name"],
                "document_hash": metadata_fields["document_hash"],
            }

            vector_records.append(record)

        logger.info(f"Prepared {len(vector_records)} records for vector storage")
        return vector_records

    async def store_document_chunks(
        self,
        chunks: List[DocumentChunk],
        namespace: str,
        vector_service: "VectorService",
    ) -> bool:
        """
        Store document chunks in vector database using batch processing.

        Args:
            chunks: List of DocumentChunk objects to store
            namespace: Pinecone namespace to store chunks in
            vector_service: VectorService instance for storage operations

        Returns:
            True if storage was successful, False otherwise
        """
        try:
            if not chunks:
                logger.warning("No chunks to store")
                return True

            # Extract document hash from namespace (format: doc_{hash})
            document_hash = (
                namespace.replace("doc_", "")
                if namespace.startswith("doc_")
                else namespace
            )

            # Prepare chunks for vector storage
            vector_records = self.prepare_chunks_for_vector_storage(
                chunks, document_hash
            )

            # Store using batch processing with 90-record chunks as per pinetest.py
            success = await vector_service.batch_upsert_records(
                namespace, vector_records
            )

            if success:
                logger.info(
                    f"Successfully stored {len(chunks)} chunks in namespace {namespace}"
                )
            else:
                logger.error(f"Failed to store chunks in namespace {namespace}")

            return success

        except Exception as e:
            logger.error(f"Failed to store document chunks: {e}")
            return False

    def _process_xlsx(self, file_path: str, document_id: str) -> List[DocumentChunk]:
        """
        Process XLSX document and extract text chunks using openpyxl.
        Converts worksheets to markdown tables, handles empty cells and multiple sheets.

        Args:
            file_path: Path to the XLSX file
            document_id: Unique identifier for the document

        Returns:
            List of DocumentChunk objects with sheet-based metadata
        """
        chunks = []
        chunk_number = 1

        try:
            # Load the workbook
            workbook = load_workbook(file_path, read_only=True, data_only=True)

            logger.info(
                f"Processing XLSX with {len(workbook.sheetnames)} worksheets: {workbook.sheetnames}"
            )

            # Process each worksheet
            for sheet_index, sheet_name in enumerate(workbook.sheetnames):
                try:
                    worksheet = workbook[sheet_name]

                    # Get all data from the worksheet
                    data_rows = []
                    max_row = worksheet.max_row
                    max_col = worksheet.max_column

                    # Skip empty worksheets
                    if max_row == 1 and max_col == 1:
                        cell_value = worksheet.cell(row=1, column=1).value
                        if cell_value is None or str(cell_value).strip() == "":
                            logger.debug(f"Skipping empty worksheet: {sheet_name}")
                            continue

                    # Extract data from all cells
                    for row in range(1, max_row + 1):
                        row_data = []
                        has_data = False

                        for col in range(1, max_col + 1):
                            cell = worksheet.cell(row=row, column=col)
                            cell_value = cell.value

                            # Handle different cell value types
                            if cell_value is None:
                                cell_text = ""
                            elif isinstance(cell_value, (int, float)):
                                cell_text = str(cell_value)
                            else:
                                cell_text = str(cell_value).strip()

                            row_data.append(cell_text)
                            if cell_text:
                                has_data = True

                        # Only add rows that have at least some data
                        if has_data:
                            data_rows.append(row_data)

                    # Convert to markdown table format if we have data
                    if data_rows:
                        markdown_content = self._convert_to_markdown_table(
                            data_rows, sheet_name
                        )

                        if markdown_content:
                            metadata = ChunkMetadata(
                                document_id=document_id,
                                page_number=sheet_index + 1,  # 1-based sheet numbering
                                line_number=1,
                                chunk_number=chunk_number,
                                document_type="xlsx",
                                file_name=Path(file_path).name,
                                sheet_name=sheet_name,
                                sheet_index=sheet_index,
                            )

                            chunk = DocumentChunk(
                                content=markdown_content, metadata=metadata
                            )
                            chunks.append(chunk)
                            chunk_number += 1

                            logger.debug(
                                f"Processed worksheet '{sheet_name}': {len(data_rows)} rows"
                            )
                    else:
                        logger.debug(f"No data found in worksheet: {sheet_name}")

                except Exception as e:
                    logger.warning(f"Failed to process worksheet '{sheet_name}': {e}")
                    continue

            workbook.close()

        except Exception as e:
            logger.error(f"Failed to process XLSX: {e}")
            raise Exception(f"Failed to process XLSX document: {e}")

        if not chunks:
            raise Exception("No readable content found in XLSX document")

        logger.info(
            f"Successfully processed XLSX: {len(chunks)} chunks from {len(workbook.sheetnames) if 'workbook' in locals() else 'unknown'} worksheets"
        )
        return chunks

    def _convert_to_markdown_table(
        self, data_rows: List[List[str]], sheet_name: str
    ) -> str:
        """
        Convert worksheet data to markdown table format.

        Args:
            data_rows: List of rows, each row is a list of cell values
            sheet_name: Name of the worksheet

        Returns:
            Markdown formatted table string
        """
        if not data_rows:
            return ""

        try:
            # Start with sheet name as header
            markdown_lines = [f"## {sheet_name}", ""]

            # Determine the maximum number of columns
            max_cols = max(len(row) for row in data_rows)

            # Pad all rows to have the same number of columns
            normalized_rows = []
            for row in data_rows:
                padded_row = row + [""] * (max_cols - len(row))
                normalized_rows.append(padded_row)

            # Create markdown table
            if normalized_rows:
                # Header row (first row)
                header_row = normalized_rows[0]
                markdown_lines.append("| " + " | ".join(header_row) + " |")

                # Separator row
                separator = "| " + " | ".join(["---"] * max_cols) + " |"
                markdown_lines.append(separator)

                # Data rows (remaining rows)
                for row in normalized_rows[1:]:
                    markdown_lines.append("| " + " | ".join(row) + " |")

            return "\n".join(markdown_lines)

        except Exception as e:
            logger.warning(
                f"Failed to convert worksheet '{sheet_name}' to markdown: {e}"
            )
            # Fallback: return raw text
            fallback_lines = [f"## {sheet_name}", ""]
            for row in data_rows:
                fallback_lines.append(" | ".join(str(cell) for cell in row))
            return "\n".join(fallback_lines)

    async def _process_image(
        self, file_path: str, document_id: str
    ) -> List[DocumentChunk]:
        """
        Process image document for direct Gemini vision API handling.
        Creates a single chunk that indicates this is an image document.

        Args:
            file_path: Path to the image file
            document_id: Unique identifier for the document

        Returns:
            List containing a single DocumentChunk with image metadata
        """
        try:
            # Create a single chunk to indicate this is an image document
            # The actual image processing will be handled directly by Gemini vision API
            metadata = ChunkMetadata(
                document_id=document_id,
                page_number=1,
                line_number=1,
                chunk_number=1,
                document_type="image",
                file_name=Path(file_path).name,
            )

            # Create a placeholder chunk - actual image analysis will be done by Gemini vision API
            chunk_content = f"Image document: {Path(file_path).name}"
            chunk = DocumentChunk(content=chunk_content, metadata=metadata)

            logger.info(f"Processed image document: {Path(file_path).name}")
            return [chunk]

        except Exception as e:
            logger.error(f"Failed to process image: {e}")
            raise Exception(f"Failed to process image document: {e}")


class VectorService:
    """
    Pinecone vector service following the patterns from tests/pinetest.py.
    Handles document storage, search, and namespace management.
    """

    def __init__(self):
        self.pc = Pinecone(api_key=Config.PINECONE_API_KEY)
        self.index = self.pc.Index(Config.PINECONE_INDEX_NAME)
        # Add semaphore for concurrent request limiting
        self._semaphore = asyncio.Semaphore(Config.MAX_CONCURRENT_REQUESTS)

        # Initialize reranker service once to avoid duplicate initializations
        self._reranker_service = None

    def get_reranker_service(self):
        """Get or initialize the reranker service (lazy initialization)."""
        if self._reranker_service is None:
            # Import RerankerService here to avoid circular imports
            from app.reranker_services import RerankerService

            self._reranker_service = RerankerService()
            logger.info("RerankerService initialized once for VectorService")
        return self._reranker_service

    async def check_namespace_exists(self, namespace: str) -> bool:
        """
        Check if a namespace exists and has vectors for hash-based caching.

        Args:
            namespace: Namespace to check (format: doc_{hash})

        Returns:
            True if namespace exists with vectors, False otherwise
        """
        try:
            logger.debug(f"Checking if namespace exists: {namespace}")

            stats = await asyncio.to_thread(self.index.describe_index_stats)
            namespaces = stats.get("namespaces", {})

            if namespace in namespaces:
                vector_count = namespaces[namespace].get("vector_count", 0)
                logger.debug(
                    f"Namespace {namespace} exists with {vector_count} vectors"
                )
                return vector_count > 0

            logger.debug(f"Namespace {namespace} does not exist")
            return False

        except Exception as e:
            logger.warning(f"Failed to check namespace {namespace}: {e}")
            return False

    async def batch_upsert_records(
        self, namespace: str, records: List[Dict[str, Any]]
    ) -> bool:
        """
        Batch upsert records to Pinecone following pinetest.py patterns.
        Uses 90-record batches as shown in pinetest.py example.

        Args:
            namespace: Pinecone namespace to upsert records to
            records: List of records in format {"_id": "unique-id", "text": "content", "metadata": "fields"}

        Returns:
            True if all batches were successful, False otherwise
        """
        try:
            if not records:
                logger.warning("No records to upsert")
                return True

            nRecords = len(records)

            logger.info(
                f"Starting batch upsert of {nRecords} records to namespace {namespace}"
            )

            # Use semaphore to limit concurrent operations
            async with self._semaphore:
                # Process records in batches - increased to 96 for better performance
                batch_size = 95
                successful_batches = 0
                total_batches = nRecords // batch_size + (
                    1 if nRecords % batch_size else 0
                )

                # Use the chunks helper function from pinetest.py
                for i, batch_records in enumerate(chunks(records, batch_size)):
                    try:
                        # Convert tuple back to list for processing
                        batch_list = list(batch_records)

                        # Upsert batch using Pinecone upsert_records method (for built-in embeddings)
                        await asyncio.to_thread(
                            self.index.upsert_records, namespace, batch_list
                        )

                        successful_batches += 1
                        logger.info(
                            f"Successfully upserted batch {i+1} with {len(batch_list)} records"
                        )

                        # Small delay between batches to avoid rate limiting
                        if total_batches > 5:
                            logger.info(
                                f"WAITING! - Large Document - batch {i+1}"
                            )
                            await asyncio.sleep(2.5)
                        else:
                            logger.info(f"WAITING! - Small Document - batch {i+1}")
                            await asyncio.sleep(0.3)  # small delay for smaller batches

                    except Exception as batch_error:
                        logger.error(f"Failed to upsert batch {i+1}: {batch_error}")
                        continue

                success_rate = (
                    successful_batches / total_batches if total_batches > 0 else 0
                )
                logger.info(
                    f"Batch upsert completed: {successful_batches}/{total_batches} batches successful ({success_rate:.1%})"
                )

                # Wait for Pinecone to index the vectors before returning
                if success_rate >= 0.8:
                    logger.info("Waiting for Pinecone to index vectors...")
                    await self._wait_for_indexing(namespace, len(records))

                # Consider successful if at least 80% of batches succeeded
                return success_rate >= 0.8

        except Exception as e:
            logger.error(f"Batch upsert failed for namespace {namespace}: {e}")
            return False

    async def _wait_for_indexing(
        self, namespace: str, expected_count: int, max_wait_time: int = None
    ):
        """
        Wait for Pinecone to finish indexing vectors after upsert.

        Args:
            namespace: Namespace to check
            expected_count: Expected number of vectors
            max_wait_time: Maximum time to wait in seconds (default: 60)
        """
        max_wait = max_wait_time or 60  # Default 60 seconds timeout
        wait_interval = 0.2
        check_count = 0
        start_time = time.time()

        while time.time() - start_time < max_wait:
            try:
                check_count += 1
                # Check namespace stats to see if vectors are indexed
                stats = await asyncio.to_thread(self.index.describe_index_stats)

                if hasattr(stats, "namespaces") and namespace in stats.namespaces:
                    namespace_stats = stats.namespaces[namespace]
                    vector_count = getattr(namespace_stats, "vector_count", 0)

                    logger.debug(
                        f"Namespace {namespace} has {vector_count} vectors indexed (expected: {expected_count}) - check #{check_count}"
                    )

                    # If we have at least 80% of expected vectors, consider it ready
                    if vector_count >= expected_count * 0.8:
                        logger.info(
                            f"Indexing complete for namespace {namespace}: {vector_count} vectors available after {check_count} checks"
                        )
                        return

                await asyncio.sleep(wait_interval)

                # Log progress every 10 checks
                if check_count % 40 == 0:
                    elapsed = time.time() - start_time
                    logger.info(
                        f"Still waiting for indexing completion for namespace {namespace} - check #{check_count} ({elapsed:.1f}s elapsed)"
                    )

            except Exception as e:
                logger.warning(f"Error checking indexing status: {e}")
                await asyncio.sleep(1.0)

        # Timeout reached
        elapsed = time.time() - start_time
        logger.warning(
            f"Indexing wait timeout reached for namespace {namespace} after {elapsed:.1f}s"
        )

    async def is_namespace_ready(self, namespace: str, min_vectors: int = 1) -> bool:
        """
        Check if a namespace is ready for search operations.

        Args:
            namespace: Namespace to check
            min_vectors: Minimum number of vectors expected

        Returns:
            True if namespace is ready, False otherwise
        """
        try:
            stats = await asyncio.to_thread(self.index.describe_index_stats)

            if hasattr(stats, "namespaces") and namespace in stats.namespaces:
                namespace_stats = stats.namespaces[namespace]
                vector_count = getattr(namespace_stats, "vector_count", 0)
                return vector_count >= min_vectors

            return False

        except Exception as e:
            logger.warning(f"Error checking namespace readiness: {e}")
            return False

    async def search_documents(
        self, query: str, namespace: str, top_k: int = None, use_semaphore: bool = True
    ) -> VectorSearchResult:
        """
        Search documents in Pinecone without reranking (raw vector search only).

        Args:
            query: Search query text
            namespace: Namespace to search in
            top_k: Number of results to return (uses Config.PINECONE_TOP_K if None)
            use_semaphore: Whether to use semaphore for rate limiting (default True)

        Returns:
            VectorSearchResult with relevant chunks
        """
        start_time = datetime.utcnow()

        # Use configured top_k if not provided
        if top_k is None:
            top_k = Config.PINECONE_TOP_K

        if use_semaphore:
            async with self._semaphore:
                return await self._search_documents_internal(query, namespace, top_k, start_time)
        else:
            return await self._search_documents_internal(query, namespace, top_k, start_time)

    async def _search_documents_internal(self, query: str, namespace: str, top_k: int, start_time: datetime) -> VectorSearchResult:
        """Internal method to search documents without semaphore."""
        try:
            logger.debug(f"Search: {query[:50]}... (top_k={top_k})")

            # Use standard semantic search without reranking
            results = await asyncio.to_thread(
                self.index.search,
                namespace=namespace,
                query={"inputs": {"text": query}, "top_k": top_k},
                fields=[
                    "text",
                    "document_id",
                    "page_number",
                    "line_number",
                    "chunk_number",
                    "document_type",
                    "file_name",
                ],
            )

            # Parse standard search results
            relevant_chunks = []
            hits = results.get("result", {}).get("hits", [])

            for hit in hits:
                hit_id = hit.get("_id", "")
                score = float(hit.get("_score", 0.0))
                fields = hit.get("fields", {})

                # Create ChunkMetadata from stored fields
                chunk_metadata = ChunkMetadata(
                    document_id=fields.get("document_id", ""),
                    page_number=int(fields.get("page_number", 0)),
                    line_number=int(fields.get("line_number", 0)),
                    chunk_number=int(fields.get("chunk_number", 0)),
                    document_type=fields.get("document_type", ""),
                    file_name=fields.get("file_name", ""),
                )

                # Create DocumentChunk
                document_chunk = DocumentChunk(
                    content=fields.get("text", ""), metadata=chunk_metadata
                )

                # Create RelevantChunk
                relevant_chunk = RelevantChunk(
                    chunk=document_chunk,
                    relevance_score=score,
                    distance=1.0 - score,
                )

                relevant_chunks.append(relevant_chunk)

            search_time = (datetime.utcnow() - start_time).total_seconds()
            logger.debug(
                f"Found {len(relevant_chunks)} chunks in {search_time:.2f}s"
            )

            return VectorSearchResult(
                chunks=relevant_chunks,
                search_time=search_time,
                total_results=len(relevant_chunks),
                )

        except Exception as e:
            logger.error(f"Search failed: {e}")
            raise Exception(f"Vector search failed: {e}")

    async def search_multiple_queries(
        self, queries: List[str], namespace: str, original_query: str, top_k: int = None
    ) -> List[RelevantChunk]:
        """
        Perform parallel vector searches for multiple queries, deduplicate, and rerank results.

        Architecture:
        1. Run each query → get top_k (17) results per query
        2. Merge all results (17 × n chunks)
        3. Remove duplicates from merged list
        4. Run unique list through reranker → get final top_n (4) chunks

        Args:
            queries: List of query strings to search
            namespace: Document namespace to search in
            original_query: Original user query for final reranking
            top_k: Number of results per query (from PINECONE_TOP_K)

        Returns:
            Final reranked list of chunks (limited to PINECONE_RERANK_TOP_N)
        """
        if not queries:
            logger.warning("No queries provided for multi-query search")
            return []

        # Use configured top_k if not provided
        if top_k is None:
            top_k = Config.PINECONE_TOP_K

        try:
            # Step 1: Run each query and get top_k results (raw vector search without semaphore)
            search_tasks = [
                self.search_documents(query, namespace, top_k, use_semaphore=False) for query in queries
            ]

            search_results = await asyncio.gather(*search_tasks, return_exceptions=True)

            # Step 2: Merge all results (17 × n chunks)
            all_chunks = []
            total_chunks_found = 0
            successful_searches = 0

            for i, result in enumerate(search_results):
                if isinstance(result, Exception):
                    logger.debug(f"Query {i+1} search failed: {result}")
                    continue

                if isinstance(result, VectorSearchResult) and result.chunks:
                    successful_searches += 1
                    chunks_found = len(result.chunks)
                    total_chunks_found += chunks_found
                    logger.debug(
                        f"Query {i+1}: found {chunks_found} chunks (expected {top_k})"
                    )

                    # Add all chunks to the merged list
                    all_chunks.extend(result.chunks)
                else:
                    logger.debug(f"Query {i+1}: no chunks found")

            # Step 3: Remove duplicates from merged list
            unique_chunks = {}  # Use dict with chunk ID as key for deduplication

            for chunk in all_chunks:
                # Generate unique chunk ID for deduplication
                chunk_id = f"{chunk.chunk.metadata.document_id}_{chunk.chunk.metadata.chunk_number}"

                # Keep the chunk with the highest relevance score if duplicate
                if (
                    chunk_id not in unique_chunks
                    or chunk.relevance_score > unique_chunks[chunk_id].relevance_score
                ):
                    unique_chunks[chunk_id] = chunk

            unique_chunk_list = list(unique_chunks.values())

            # Calculate deduplication efficiency
            deduplication_efficiency = (
                len(unique_chunk_list) / max(total_chunks_found, 1) * 100
            )
            duplicates_removed = total_chunks_found - len(unique_chunk_list)

            logger.info(
                f"🔍 Multi-query search: {successful_searches}/{len(queries)} successful, {total_chunks_found} total → {len(unique_chunk_list)} unique chunks ({duplicates_removed} duplicates removed, {deduplication_efficiency:.1f}% efficiency)"
            )

            # Step 4: Run unique list through reranker to get final top_n chunks
            if unique_chunk_list:
                final_chunks = await self.rerank_chunks_with_original_query(
                    unique_chunk_list, original_query
                )
                logger.info(
                    f"🎯 Final reranked result: {len(final_chunks)} chunks for context"
                )
                return final_chunks
            else:
                return []

        except Exception as e:
            logger.error(f"Multi-query search failed: {e}")
            return []

    async def rerank_chunks_with_original_query(
        self, chunks: List[RelevantChunk], original_query: str
    ) -> List[RelevantChunk]:
        """
        Rerank chunks using the unified RerankerService with the original user query.

        Args:
            chunks: List of chunks to rerank
            original_query: Original user question for reranking

        Returns:
            Reranked list of chunks (top N based on config)
        """
        if not chunks:
            logger.warning("No chunks provided for reranking")
            return []

        if not Config.PINECONE_RERANKING_ENABLED:
            # Return top N chunks without reranking
            logger.info(
                f"Reranking disabled, returning top {Config.PINECONE_RERANK_TOP_N} chunks"
            )
            return chunks[: Config.PINECONE_RERANK_TOP_N]

        try:
            logger.info(
                f"Reranking {len(chunks)} chunks with original query: {original_query[:50]}..."
            )

            # Use singleton reranker service to avoid duplicate initializations
            reranker_service = self.get_reranker_service()

            # Check if reranker service is available
            if not reranker_service.is_available():
                logger.warning(
                    "No reranker services available, returning top chunks without reranking"
                )
                return chunks[: Config.PINECONE_RERANK_TOP_N]

            # Prepare documents for reranking
            documents = []
            for chunk in chunks:
                documents.append(chunk.chunk.content)

            # Use unified reranker service
            reranker_results = await reranker_service.rerank(
                query=original_query,
                documents=documents,
                top_n=Config.PINECONE_RERANK_TOP_N,
            )

            # Map reranked results back to chunks
            reranked_chunks = []
            for result in reranker_results:
                # Find original chunk by index
                if 0 <= result.index < len(chunks):
                    original_chunk = chunks[result.index]

                    # Update relevance score with reranked score
                    reranked_chunk = RelevantChunk(
                        chunk=original_chunk.chunk,
                        relevance_score=result.score,
                        distance=1.0 - result.score,
                    )
                    reranked_chunks.append(reranked_chunk)

            logger.info(
                f"Reranking completed using {reranker_results[0].service_used if reranker_results else 'unknown'} service, returned {len(reranked_chunks)} chunks"
            )
            return reranked_chunks

        except Exception as e:
            logger.error(
                f"Reranking failed: {e}, returning top chunks without reranking"
            )
            return chunks[: Config.PINECONE_RERANK_TOP_N]


class LLMService:
    """
    Enhanced LLM service with multi-key support and Groq integration.
    Supports automatic failover between Gemini and Groq services.
    """

    def __init__(self):
        """Initialize LLM service with multi-key managers and service priority."""
        # Load API keys from config
        gemini_keys, groq_keys = Config.load_api_keys()
        primary_service, secondary_service = Config.get_llm_service_priority()

        # Initialize key managers
        self.gemini_key_manager = (
            LLMKeyManager("gemini", gemini_keys) if gemini_keys else None
        )
        self.groq_key_manager = LLMKeyManager("groq", groq_keys) if groq_keys else None

        # Set service priority
        self.primary_service = primary_service
        self.secondary_service = secondary_service

        # Initialize Groq client (will be set per request with different keys)
        self.groq_client = None

        # Semaphore for concurrent request limiting
        self._semaphore = asyncio.Semaphore(Config.MAX_CONCURRENT_REQUESTS)

        # Prompt templates are now loaded from centralized JSON file via Config

        logger.info(
            f"LLMService initialized - Primary: {primary_service}, Secondary: {secondary_service}"
        )
        logger.info(
            f"Gemini keys: {len(gemini_keys) if gemini_keys else 0}, Groq keys: {len(groq_keys) if groq_keys else 0}"
        )
        logger.info("Prompt templates loaded from centralized prompts.json")
        logger.info(
            "Flight number tool available ONLY for direct document upload (removed from regular text generation and images)"
        )

    async def get_flight_num(
        self, pdf_content: bytes = None, pdf_content_type: str = None
    ) -> str:
        """Gets the flight number using the provided PDF content

        Args:
            pdf_content: Optional PDF content bytes from the API request
            pdf_content_type: Optional PDF content type

        Returns:
            A string containing the flight number.
        """
        try:
            logger.info("Flight number tool called - using single LLM call approach")

            # Step 1: Get city from the API
            async with aiohttp.ClientSession() as session:
                async with session.get(
                    "https://register.hackrx.in/submissions/myFavouriteCity"
                ) as response:
                    if response.status != 200:
                        raise Exception(
                            f"Failed to fetch city data: HTTP {response.status}"
                        )
                    data = await response.json()
                    city = data["data"]["city"]
                    logger.info(f"Retrieved city: {city}")

            logger.info(f"Using provided PDF content ({len(pdf_content)} bytes)")

            prompt = Config.get_prompt("flight_pdf_analysis_prompt", city=city)

            # Make single LLM call with PDF upload
            logger.info(
                "Making SINGLE LLM call with PDF upload to get both landmark and API URL"
            )

            response_text = await self._make_single_pdf_llm_request(prompt, pdf_content, pdf_content_type)

            # Parse the JSON response
            import json

            try:
                # Strip markdown code blocks if present
                clean_response = response_text.strip()
                if clean_response.startswith("```json"):
                    clean_response = clean_response[7:]  # Remove ```json
                if clean_response.startswith("```"):
                    clean_response = clean_response[3:]  # Remove ```
                if clean_response.endswith("```"):
                    clean_response = clean_response[:-3]  # Remove trailing ```
                clean_response = clean_response.strip()

                logger.debug(f"Cleaned JSON response: {clean_response}")
                response_data = json.loads(clean_response)
                landmark = response_data.get("landmark", "").strip()
                api_url = response_data.get("api_url", "").strip()

                logger.info(f"Retrieved landmark: {landmark}")
                logger.info(f"Retrieved API URL: {api_url}")

                # Validate the results
                if not landmark:
                    raise Exception("No landmark found in LLM response")
                if not api_url or not api_url.startswith("https://"):
                    raise Exception(f"Invalid API URL: {api_url}")
                if "example.com" in api_url:
                    raise Exception(
                        f"LLM returned example URL instead of actual URL: {api_url}"
                    )

            except json.JSONDecodeError as e:
                logger.error(f"Failed to parse JSON response: {response_text}")
                raise Exception(f"LLM returned invalid JSON: {e}")

            # Step 4: Get flight number from the API endpoint
            async with aiohttp.ClientSession() as session:
                async with session.get(api_url) as response:
                    if response.status != 200:
                        raise Exception(
                            f"Failed to fetch flight data: HTTP {response.status}"
                        )
                    flight_data = await response.json()
                    flight_number = flight_data["data"]["flightNumber"]
                    logger.info(f"Retrieved flight number: {flight_number}")

            return flight_number

        except Exception as e:
            logger.error(f"Flight number tool failed: {e}")
            return f"Error retrieving flight number: {str(e)}"

    async def generate_answer(
        self, question: str, context_chunks: List[RelevantChunk], use_semaphore: bool = True
    ) -> str:
        """
        Generate answer to a question using context chunks with multi-service fallback.

        Args:
            question: User question
            context_chunks: List of relevant document chunks
            use_semaphore: Whether to use semaphore for rate limiting (default True)

        Returns:
            Generated answer string
        """
        if use_semaphore:
            async with self._semaphore:
                return await self._generate_answer_internal(question, context_chunks)
        else:
            return await self._generate_answer_internal(question, context_chunks)

    async def _generate_answer_internal(self, question: str, context_chunks: List[RelevantChunk]) -> str:
        """Internal method to generate answer without semaphore."""
        try:
            # Prepare context from chunks
            context_texts = []
            for chunk in context_chunks:
                context_texts.append(chunk.chunk.content)

            context = (
                "\n\n".join(context_texts)
                if context_texts
                else "No relevant context found."
            )

            # Get prompt template from centralized prompts.json
            prompt = Config.get_prompt(
                "decision_prompt", context=context, question=question
            )

            logger.info(f"Generating answer for question: {question[:50]}...")
            logger.debug("Using decision_prompt from centralized prompts.json")

            # Try primary service first, then fallback
            answer = await self._make_request_with_fallback(prompt)

            logger.info(f"Generated answer: {answer[:100]}...")
            return answer

        except Exception as e:
            logger.error(f"Answer generation failed: {e}")
            return f"I apologize, but I encountered an error while processing your question: {str(e)}"

    async def _make_request_with_fallback(self, prompt: str) -> str:
        """
        Make LLM request with cycling through all Gemini keys first, then Groq fallback.
        Always tries Gemini first, uses Groq only when individual Gemini requests fail.

        Args:
            prompt: Formatted prompt string

        Returns:
            Generated response string
        """
        # Always try Gemini first (cycling through all keys)
        if self.gemini_key_manager and self.gemini_key_manager.api_keys:
            try:
                return await self._try_gemini_request_with_cycling(prompt)
            except Exception as e:
                # If all Gemini keys fail, try Groq as fallback
                if self.groq_key_manager and self.groq_key_manager.api_keys:
                    try:
                        return await self._try_groq_request_with_cycling(prompt)
                    except Exception as groq_error:
                        raise Exception(
                            f"Both Gemini and Groq services failed. Gemini: {e}, Groq: {groq_error}"
                        )
                else:
                    raise Exception(f"Gemini failed and no Groq keys available: {e}")
        elif self.groq_key_manager and self.groq_key_manager.api_keys:
            # Only Groq available
            return await self._try_groq_request_with_cycling(prompt)
        else:
            raise Exception("No API keys configured for any service")

    async def _try_gemini_request_with_cycling(self, prompt: str) -> str:
        """
        Try making request to Gemini cycling through all keys.
        For each individual request, if Gemini fails, immediately try Groq.
        Keep cycling through all Gemini keys indefinitely.

        Args:
            prompt: Formatted prompt string

        Returns:
            Generated response string
        """
        if not self.gemini_key_manager:
            raise Exception("No Gemini API keys configured")

        # Try all Gemini keys once
        for attempt in range(len(self.gemini_key_manager.api_keys)):
            try:
                # Get next key in rotation
                api_key = await self.gemini_key_manager.get_next_key()

                # Configure Gemini with current key
                genai.configure(api_key=api_key)

                # Create model instance WITHOUT flight tool for regular text generation
                # Flight tool should only be available when processing flight-related documents
                model = genai.GenerativeModel(
                    model_name=Config.GEMINI_MODEL,
                    generation_config={
                        "temperature": Config.GEMINI_TEMPERATURE,
                        "top_p": Config.GEMINI_TOP_P,
                        "max_output_tokens": Config.GEMINI_MAX_TOKENS,
                    },
                    safety_settings={
                        HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                        HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                        HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                        HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                    },
                    # NO TOOLS - flight tool only available for flight-related documents
                )

                # Make request
                response = await asyncio.to_thread(model.generate_content, prompt)

                if not response:
                    raise Exception("Empty response from Gemini API")

                logger.debug(f"Gemini response type: {type(response)}")
                logger.debug(
                    f"Gemini response candidates: {len(response.candidates) if response.candidates else 0}"
                )

                # Return regular text response (no function calls in regular text generation)
                if response.text:
                    # Add sleep after successful LLM call
                    await asyncio.sleep(0.5)
                    return response.text.strip()
                else:
                    raise Exception("No text response from Gemini API")

            except Exception as e:
                # Add sleep after failed LLM call
                await asyncio.sleep(0.5)

                # For individual Gemini failure, try Groq immediately
                if self.groq_key_manager and self.groq_key_manager.api_keys:
                    try:
                        return await self._try_single_groq_request(prompt)
                    except Exception:
                        # Continue to next Gemini key if Groq also fails
                        pass

                # Continue to next Gemini key
                continue

        # If all Gemini keys failed, raise exception
        raise Exception("All Gemini keys exhausted")

    async def _try_gemini_request_without_tools(self, prompt: str) -> str:
        """
        Try making request to Gemini WITHOUT flight tool (for internal use).

        Args:
            prompt: Formatted prompt string

        Returns:
            Generated response string
        """
        if not self.gemini_key_manager:
            raise Exception("No Gemini API keys configured")

        # Try all Gemini keys once
        for attempt in range(len(self.gemini_key_manager.api_keys)):
            try:
                # Get next key in rotation
                api_key = await self.gemini_key_manager.get_next_key()

                # Configure Gemini with current key
                genai.configure(api_key=api_key)

                # Create model instance WITHOUT flight number tool
                model = genai.GenerativeModel(
                    model_name=Config.GEMINI_MODEL,
                    generation_config={
                        "temperature": Config.GEMINI_TEMPERATURE,
                        "top_p": Config.GEMINI_TOP_P,
                        "max_output_tokens": Config.GEMINI_MAX_TOKENS,
                    },
                    safety_settings={
                        HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                        HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                        HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                        HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
                    },
                    # No tools parameter - this is the key difference
                )

                # Make request
                response = await asyncio.to_thread(model.generate_content, prompt)

                if not response or not response.text:
                    raise Exception("Empty response from Gemini API")

                # Add sleep after successful LLM call
                await asyncio.sleep(0.5)

                return response.text.strip()

            except Exception as e:
                # Add sleep after failed LLM call
                await asyncio.sleep(0.5)

                # For individual Gemini failure, try Groq immediately
                if self.groq_key_manager and self.groq_key_manager.api_keys:
                    try:
                        return await self._try_single_groq_request(prompt)
                    except Exception:
                        # Continue to next Gemini key if Groq also fails
                        pass

                # Continue to next Gemini key
                continue

        # If all Gemini keys failed, raise exception
        raise Exception("All Gemini keys exhausted (internal request)")

    async def _make_single_pdf_llm_request(
        self, prompt: str, pdf_content: bytes, content_type: str
    ) -> str:
        """
        Make a single LLM request with PDF upload (without flight tool to avoid recursion).
        Used specifically for the flight tool's internal PDF analysis.

        Args:
            prompt: Text prompt for PDF analysis
            pdf_content: PDF content as bytes
            content_type: Content type of the PDF

        Returns:
            Generated response string
        """
        if not self.gemini_key_manager:
            raise Exception("No Gemini API keys configured")

        # Try all Gemini keys
        for attempt in range(len(self.gemini_key_manager.api_keys)):
            try:
                # Get next key in rotation
                api_key = await self.gemini_key_manager.get_next_key()

                # Use the new Google GenAI client approach
                # Create client with current key
                client = google_genai.Client(api_key=api_key)

                # Create document part from bytes
                document_part = types.Part.from_bytes(
                    data=pdf_content, mime_type=content_type
                )

                # Make request with document and prompt (NO TOOLS to avoid recursion)
                response = await asyncio.to_thread(
                    client.models.generate_content,
                    model="gemini-2.5-flash",
                    contents=[prompt, document_part],
                    config=types.GenerateContentConfig(
                        temperature=Config.GEMINI_TEMPERATURE,
                        top_p=Config.GEMINI_TOP_P,
                        max_output_tokens=Config.GEMINI_MAX_TOKENS,
                        thinking_config=types.ThinkingConfig(thinking_budget=0),
                        # NO TOOLS - this is crucial to avoid recursion
                    ),
                )

                if not response or not response.text:
                    raise Exception("Empty response from Gemini document API")

                # Add sleep after successful LLM call
                await asyncio.sleep(0.5)

                logger.info(f"Successfully processed PDF with single Gemini call")
                return response.text.strip()

            except Exception as e:
                # Add sleep after failed LLM call
                await asyncio.sleep(0.5)

                # Continue to next Gemini key
                logger.warning(
                    f"Single PDF LLM request failed with key {attempt + 1}: {e}"
                )
                continue

        # If all Gemini keys failed, raise exception
        raise Exception("All Gemini keys exhausted for single PDF request")

    async def _try_single_groq_request(self, prompt: str) -> str:
        """
        Try making a single request to Groq with the next key.

        Args:
            prompt: Formatted prompt string

        Returns:
            Generated response string
        """
        if not self.groq_key_manager:
            raise Exception("No Groq API keys configured")

        # Get next Groq key
        api_key = await self.groq_key_manager.get_next_key()

        # Create Groq client with current key
        client = Groq(api_key=api_key)

        # Make request
        completion = await asyncio.to_thread(
            client.chat.completions.create,
            model=Config.GROQ_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": "You are an insurance policy expert who provides clear, concise answers in simple language.",
                },
                {"role": "user", "content": prompt},
            ],
            temperature=Config.GROQ_TEMPERATURE,
            max_tokens=Config.GROQ_MAX_TOKENS,
            top_p=Config.GROQ_TOP_P,
            stream=False,
            stop=None,
        )

        if not completion.choices or not completion.choices[0].message.content:
            raise Exception("Empty response from Groq API")

        # Add sleep after LLM call
        await asyncio.sleep(0.5)

        return completion.choices[0].message.content.strip()

    async def _try_groq_request_with_cycling(self, prompt: str) -> str:
        """
        Try making request to Groq cycling through all keys.

        Args:
            prompt: Formatted prompt string

        Returns:
            Generated response string
        """
        if not self.groq_key_manager:
            raise Exception("No Groq API keys configured")

        # Try all Groq keys once
        for attempt in range(len(self.groq_key_manager.api_keys)):
            try:
                return await self._try_single_groq_request(prompt)
            except Exception as e:
                # Continue to next Groq key
                continue

        # If all Groq keys failed, raise exception
        raise Exception("All Groq keys exhausted")

    async def generate_image_answer(self, question: str, image_path: str) -> str:
        """
        Generate answer to a question about an image using Gemini vision API.
        Bypasses vector database and relevance checking for direct image analysis.

        Args:
            question: User question about the image
            image_path: Path to the image file

        Returns:
            Generated answer string based on image analysis
        """
        async with self._semaphore:
            try:
                logger.info(f"Generating image answer for question: {question[:50]}...")

                # Read image file
                with open(image_path, "rb") as image_file:
                    image_data = image_file.read()

                # Create image part for Gemini vision API
                import google.generativeai as genai
                from PIL import Image
                import io

                # Load image using PIL
                image = Image.open(io.BytesIO(image_data))

                # Create prompt for image analysis
                prompt = Config.get_prompt("decision_prompt", question=question, context="The attached image is the provided context.")

                # Try Gemini first with vision capabilities
                if self.gemini_key_manager and self.gemini_key_manager.api_keys:
                    try:
                        return await self._try_gemini_vision_request(prompt, image)
                    except Exception as e:
                        logger.error(f"Gemini vision API failed: {e}")
                        return f"I apologize, but I encountered an error while analyzing the image: {str(e)}"
                else:
                    return "I apologize, but image analysis is not available as no Gemini API keys are configured."

            except Exception as e:
                logger.error(f"Image answer generation failed: {e}")
                return f"I apologize, but I encountered an error while processing the image: {str(e)}"

    async def _try_gemini_vision_request(self, prompt: str, image) -> str:
        """
        Try making a vision request to Gemini with image analysis.

        Args:
            prompt: Text prompt for image analysis
            image: PIL Image object

        Returns:
            Generated response string
        """
        if not self.gemini_key_manager:
            raise Exception("No Gemini API keys configured")

        # Try all Gemini keys
        for attempt in range(len(self.gemini_key_manager.api_keys)):
            try:
                # Get next key in rotation
                api_key = await self.gemini_key_manager.get_next_key()

                # Use the new Google GenAI client approach for better control
                client = google_genai.Client(api_key=api_key)

                # Create image part from PIL Image
                import io

                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format="PNG")
                img_byte_arr = img_byte_arr.getvalue()

                image_part = types.Part.from_bytes(
                    data=img_byte_arr, mime_type="image/png"
                )

                # Make request with image and prompt, thinking budget set to 0
                response = await asyncio.to_thread(
                    client.models.generate_content,
                    model="gemini-2.5-flash",
                    contents=[prompt, image_part],
                    config=types.GenerateContentConfig(
                        temperature=Config.GEMINI_TEMPERATURE,
                        top_p=Config.GEMINI_TOP_P,
                        max_output_tokens=Config.GEMINI_MAX_TOKENS,
                        thinking_config=types.ThinkingConfig(thinking_budget=0),
                        # No tools for image analysis
                    ),
                )

                if not response or not response.text:
                    raise Exception("Empty response from Gemini vision API")

                # Add sleep after successful LLM call
                await asyncio.sleep(0.5)

                return response.text.strip()

            except Exception as e:
                # Add sleep after failed LLM call
                await asyncio.sleep(0.5)

                # Continue to next Gemini key
                logger.warning(
                    f"Gemini vision request failed with key {attempt + 1}: {e}"
                )
                continue

        # If all Gemini keys failed, raise exception
        raise Exception("All Gemini vision API keys exhausted")

    async def generate_direct_gemini_answer(
        self, question: str, content: bytes, content_type: str, filename: str
    ) -> str:
        """
        Generate answer to a question about a document using direct Gemini upload.
        Bypasses vector database and uploads the document directly to Gemini.

        Args:
            question: User question about the document
            content: Raw document content as bytes
            content_type: MIME type of the document
            filename: Name of the document file

        Returns:
            Generated answer string based on direct document analysis
        """
        async with self._semaphore:
            try:
                logger.info(
                    f"Generating direct Gemini answer for question: {question[:50]}... on file: {filename}"
                )

                # Determine MIME type for Gemini
                mime_type = self._get_gemini_mime_type(content_type, filename)
                logger.debug(
                    f"Content type mapping: '{content_type}' -> '{mime_type}' for file: {filename}"
                )

                # Create prompt for document analysis
                prompt = Config.get_prompt(
                    "direct_document_analysis_prompt", question=question
                )

                # Try Gemini with document upload capabilities
                if self.gemini_key_manager and self.gemini_key_manager.api_keys:
                    try:
                        return await self._try_gemini_document_request(
                            prompt, content, mime_type, filename
                        )
                    except Exception as e:
                        logger.error(f"Gemini document API failed: {e}")
                        return f"I apologize, but I encountered an error while analyzing the document: {str(e)}"
                else:
                    return "I apologize, but document analysis is not available as no Gemini API keys are configured."

            except Exception as e:
                logger.error(f"Direct Gemini answer generation failed: {e}")
                return f"I apologize, but I encountered an error while processing the document: {str(e)}"

    def _get_gemini_mime_type(self, content_type: str, filename: str) -> str:
        """
        Get the appropriate MIME type for Gemini API based on content type and filename.

        Args:
            content_type: Original content type
            filename: Document filename

        Returns:
            MIME type suitable for Gemini API
        """
        # Map common content types to Gemini-supported MIME types
        mime_mapping = {
            "application/pdf": "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "text/plain": "text/plain",
            "text/html": "text/plain",  # Convert HTML to text/plain for Gemini
            "text/html; charset=utf-8": "text/plain",  # Handle charset variants
            "text/html; charset=UTF-8": "text/plain",
            "application/json": "text/plain",  # JSON as text
            "application/xml": "text/plain",  # XML as text
            "text/xml": "text/plain",
            "text/css": "text/plain",
            "text/javascript": "text/plain",
            "application/javascript": "text/plain",
            "image/jpeg": "image/jpeg",
            "image/png": "image/png",
            "image/gif": "image/gif",
            "image/webp": "image/webp",
        }

        # Handle content types with charset or other parameters
        base_content_type = content_type.split(";")[0].strip() if content_type else ""

        # Try exact match first
        if content_type in mime_mapping:
            return mime_mapping[content_type]

        # Try base content type (without charset)
        if base_content_type in mime_mapping:
            return mime_mapping[base_content_type]

        # Handle text/* types generically
        if base_content_type.startswith("text/"):
            return "text/plain"

        # Fallback to file extension
        extension = Path(filename).suffix.lower()
        extension_mapping = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".txt": "text/plain",
            ".html": "text/plain",
            ".htm": "text/plain",
            ".xml": "text/plain",
            ".json": "text/plain",
            ".css": "text/plain",
            ".js": "text/plain",
            ".md": "text/plain",
            ".csv": "text/plain",
            ".log": "text/plain",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }

        # If we have an extension match, use it
        if extension in extension_mapping:
            return extension_mapping[extension]

        # Final fallback: if it looks like a web URL or has no extension, assume text
        if (
            not extension
            or filename.startswith("http")
            or "get-secret-token" in filename
        ):
            return "text/plain"

        # Last resort fallback
        return "text/plain"  # Changed from application/octet-stream to text/plain

    async def _try_gemini_document_request(
        self, prompt: str, content: bytes, mime_type: str, filename: str
    ) -> str:
        """
        Try making a document request to Gemini with direct document upload.

        Args:
            prompt: Text prompt for document analysis
            content: Document content as bytes
            mime_type: MIME type of the document
            filename: Document filename

        Returns:
            Generated response string
        """
        if not self.gemini_key_manager:
            raise Exception("No Gemini API keys configured")

        # Try all Gemini keys
        for attempt in range(len(self.gemini_key_manager.api_keys)):
            try:
                # Get next key in rotation
                api_key = await self.gemini_key_manager.get_next_key()

                # Use the new Google GenAI client approach
                # Create client with current key
                client = google_genai.Client(api_key=api_key)

                # Create document part from bytes
                document_part = types.Part.from_bytes(data=content, mime_type=mime_type)

                # Create flight tool for new GenAI client
                flight_tool_func = types.Tool(
                    function_declarations=[
                        types.FunctionDeclaration(
                            name="get_flight_number",
                            description="Gets the user's flight number by looking up their favorite city and finding the corresponding landmark, then retrieving the flight number from the appropriate API endpoint. Use this function whenever the user asks 'What is my flight number?' or any variation of asking for their flight number. DO NOT USE IT IN ANY OTHER INSTANCES.",
                            parameters=types.Schema(
                                type=types.Type.OBJECT, properties={}, required=[]
                            ),
                        )
                    ]
                )

                # Make request with document, prompt, and flight tool
                response = await asyncio.to_thread(
                    client.models.generate_content,
                    model=Config.GEMINI_MODEL,
                    contents=[prompt, document_part],
                    config=types.GenerateContentConfig(
                        temperature=Config.GEMINI_TEMPERATURE,
                        top_p=Config.GEMINI_TOP_P,
                        max_output_tokens=Config.GEMINI_MAX_TOKENS,
                        thinking_config=types.ThinkingConfig(thinking_budget=0),
                        tools=[flight_tool_func],
                    ),
                )

                # Check if Gemini wants to call the flight function
                if response.candidates and response.candidates[0].content.parts:
                    for part in response.candidates[0].content.parts:
                        if hasattr(part, "function_call") and part.function_call:
                            if part.function_call.name == "get_flight_number":
                                logger.info(
                                    "Gemini requested flight number tool execution in direct document mode"
                                )

                                # Execute the flight number tool with PDF content
                                flight_number = await self.get_flight_num(
                                    content, mime_type
                                )

                                # Return the flight number directly with a user-friendly message
                                return f"Your flight number is: {flight_number}"

                if not response or not response.text:
                    raise Exception("Empty response from Gemini document API")

                logger.info(
                    f"Successfully processed document {filename} with Gemini direct upload"
                )
                return response.text.strip()

            except Exception as e:

                # Continue to next Gemini key
                logger.warning(
                    f"Gemini document request failed with key {attempt + 1}: {e}"
                )
                continue

        # If all Gemini keys failed, raise exception
        raise Exception("All Gemini document API keys exhausted")


class RelevanceService:
    """
    Service for checking query relevance using document-based similarity.
    Uses ONLY vector similarity search - NO LLM calls to save API costs.
    Filters irrelevant queries before they reach expensive LLM operations.
    """

    def __init__(self):
        self.pc = Pinecone(api_key=Config.PINECONE_API_KEY)
        self.index = self.pc.Index(Config.PINECONE_INDEX_NAME)
        self.document_relevance_threshold = Config.RELEVANCE_THRESHOLD
        self.relevance_check_enabled = Config.RELEVANCE_CHECK_ENABLED
        # Add semaphore for concurrent request limiting
        self._semaphore = asyncio.Semaphore(Config.MAX_CONCURRENT_REQUESTS)

    async def check_query_relevance(
        self, query: str, query_num: int, document_namespace: str = None, use_semaphore: bool = True
    ) -> bool:
        """
        Check if a query is relevant using document-based similarity.

        Args:
            query: Query text to check for relevance
            document_namespace: Document namespace for document-based relevance checking
            use_semaphore: Whether to use semaphore for rate limiting (default True)

        Returns:
            True if query is relevant, False otherwise
        """
        if use_semaphore:
            async with self._semaphore:
                return await self._check_relevance_internal(query, query_num, document_namespace)
        else:
            return await self._check_relevance_internal(query, query_num, document_namespace)

    async def _check_relevance_internal(self, query: str, query_num: int, document_namespace: str = None) -> bool:
        """Internal method to check relevance without semaphore."""
        try:
            # Check toggle first
            if not self.relevance_check_enabled:
                logger.info("Relevance checking disabled, bypassing relevance check")
                return True

            # Require document namespace for relevance checking
            if not document_namespace:
                logger.warning(
                    "No document namespace provided for relevance check, defaulting to relevant"
                )
                return True

            # Use document-based relevance
            return await self._check_document_relevance(
                query, query_num, document_namespace
            )
        except Exception as e:
            logger.error(f"Query relevance check failed: {e}")
            # Default to relevant on error to avoid blocking valid queries
            return True

    async def _check_document_relevance(
        self, query: str, query_num: int, document_namespace: str
    ) -> bool:
        """
        Check if query is relevant to document content using document-based similarity.

        Args:
            query: Query text to check for relevance
            document_namespace: Document namespace to search against

        Returns:
            True if query is relevant to document content, False otherwise
        """
        # Use semaphore to limit concurrent requests
        async with self._semaphore:
            try:
                logger.debug(f"Relevance check: {query[:50]}...")

                # Use Pinecone semantic search against document chunks
                results = await asyncio.to_thread(
                    self.index.search,
                    namespace=document_namespace,
                    query={
                        "inputs": {"text": query},
                        "top_k": min(
                            Config.PINECONE_TOP_K, 5
                        ),  # Use smaller number for relevance checking
                    },
                    fields=["content", "metadata"],
                )

                # Parse results from Pinecone search response
                hits = results.get("result", {}).get("hits", [])

                if not hits:
                    logger.info(
                        "No document chunks found for relevance check, defaulting to relevant"
                    )
                    return True

                # Get highest similarity score
                max_score = max(hit.get("_score", 0.0) for hit in hits)
                is_relevant = max_score >= self.document_relevance_threshold

                # Log relevance result
                if is_relevant:
                    logger.debug(
                        f"Q{query_num} Relevant (score: {max_score:.3f}, threshold: {self.document_relevance_threshold})"
                    )
                else:
                    logger.info(
                        f"Q{query_num} Not relevant (score: {max_score:.3f}, threshold: {self.document_relevance_threshold})"
                    )

                return is_relevant

            except Exception as e:
                logger.error(f"Document-based relevance check failed: {e}")
                # Default to relevant on error to avoid blocking valid queries
                return True


class QueryExpansionService:
    """
    Service for generating multiple diverse queries from a single user question
    using Groq API with Gemini fallback.
    """

    def __init__(self):
        """Initialize QueryExpansionService with LLM key managers."""
        # Load API keys from config
        gemini_keys, groq_keys = Config.load_api_keys()

        # Initialize key managers
        self.groq_key_manager = LLMKeyManager("groq", groq_keys) if groq_keys else None
        self.gemini_key_manager = (
            LLMKeyManager("gemini", gemini_keys) if gemini_keys else None
        )

        # Rephrase prompt is now loaded from centralized JSON file via Config

        # Semaphore for concurrent request limiting
        self._semaphore = asyncio.Semaphore(Config.MAX_CONCURRENT_REQUESTS)

        logger.info(
            f"QueryExpansionService initialized with Groq keys: {len(groq_keys) if groq_keys else 0}, Gemini keys: {len(gemini_keys) if gemini_keys else 0}"
        )

    async def generate_diverse_queries(self, original_query: str, use_semaphore: bool = True) -> List[str]:
        """
        Generate 3 diverse queries from the original query.

        Args:
            original_query: The user's original question
            use_semaphore: Whether to use semaphore for rate limiting (default True)

        Returns:
            List of 3 diverse queries (or just original if generation fails)
        """
        if use_semaphore:
            async with self._semaphore:
                return await self._generate_queries_internal(original_query)
        else:
            return await self._generate_queries_internal(original_query)

    async def _generate_queries_internal(self, original_query: str) -> List[str]:
        """Internal method to generate queries without semaphore."""
        try:
            # Try Groq first
            if self.groq_key_manager and self.groq_key_manager.api_keys:
                try:
                    queries = await self._generate_with_groq(original_query)
                    if queries and len(queries) >= 3:
                        self._log_generated_queries(
                            original_query, queries[:3], "Groq"
                        )
                        return queries[:3]
                except Exception as e:
                    logger.debug(f"Groq query generation failed: {e}")

            # Fallback to Gemini
            if self.gemini_key_manager and self.gemini_key_manager.api_keys:
                try:
                    queries = await self._generate_with_gemini(original_query)
                    if queries and len(queries) >= 3:
                        self._log_generated_queries(
                            original_query, queries[:3], "Gemini"
                        )
                        return queries[:3]
                except Exception as e:
                    logger.debug(f"Gemini query generation failed: {e}")

            # Final fallback - return original query
            self._log_fallback_query(original_query)
            return [original_query]

        except Exception as e:
            logger.error(f"❌ Query expansion failed: {e}")
            self._log_fallback_query(original_query)
            return [original_query]

    def _log_generated_queries(
        self, original_query: str, generated_queries: List[str], method: str
    ):
        """
        Log the generated queries in a structured format for easy monitoring.

        Args:
            original_query: The original user question
            generated_queries: List of generated diverse queries
            method: The method used (Groq/Gemini)
        """
        logger.info(
            f"🎯 Multi-query generated ({method}) for {original_query}: {len(generated_queries)} queries: {generated_queries}"
        )

    def _log_fallback_query(self, original_query: str):
        """
        Log when falling back to original query only.

        Args:
            original_query: The original user question
        """
        logger.info(f"🔄 Multi-query fallback: using original query only")
        logger.info(
            f'QUERY_EXPANSION: original="{original_query}" method="fallback" count=1'
        )

    async def _generate_with_groq(self, query: str) -> List[str]:
        """
        Generate queries using Groq API.

        Args:
            query: Original query string

        Returns:
            List of generated queries
        """
        try:
            # Get API key
            api_key = await self.groq_key_manager.get_next_key()

            # Initialize Groq client
            groq_client = Groq(api_key=api_key)

            # Format prompt
            prompt = Config.get_prompt("rephrase_prompt", question=query)

            # Make request to Groq using asyncio.to_thread for true async
            completion = await asyncio.to_thread(
                groq_client.chat.completions.create,
                model=Config.GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=Config.GROQ_TEMPERATURE,
                max_tokens=Config.GROQ_MAX_TOKENS,
                top_p=Config.GROQ_TOP_P,
            )

            response_text = completion.choices[0].message.content.strip()



            # Parse response into individual queries
            queries = self._parse_query_response(response_text)

            logger.debug(f"Groq generated {len(queries)} queries")
            return queries

        except Exception as e:
            logger.error(f"Groq query generation failed: {e}")
            raise

    async def _generate_with_groq_direct(self, query: str, api_key: str) -> List[str]:
        """
        Generate queries using Groq API with pre-assigned key (bypasses key manager).

        Args:
            query: Original query string
            api_key: Pre-assigned API key

        Returns:
            List of generated queries
        """
        try:
            if not api_key:
                # Fallback to original method if no key provided
                return await self._generate_queries_internal(query)

            # Initialize Groq client with pre-assigned key
            groq_client = Groq(api_key=api_key)

            # Format prompt
            prompt = Config.get_prompt("rephrase_prompt", question=query)

            # Make request to Groq using asyncio.to_thread for true async
            completion = await asyncio.to_thread(
                groq_client.chat.completions.create,
                model=Config.GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=Config.GROQ_TEMPERATURE,
                max_tokens=Config.GROQ_MAX_TOKENS,
                top_p=Config.GROQ_TOP_P,
            )

            response_text = completion.choices[0].message.content.strip()

            # Parse response into individual queries
            queries = self._parse_query_response(response_text)

            if queries and len(queries) >= 3:
                self._log_generated_queries(query, queries[:3], "Groq")
                return queries[:3]
            else:
                # Fallback to original query
                self._log_fallback_query(query)
                return [query]

        except Exception as e:
            logger.error(f"Groq direct query generation failed: {e}")
            # Fallback to original query
            self._log_fallback_query(query)
            return [query]

    async def _generate_with_gemini(self, query: str) -> List[str]:
        """
        Generate queries using Gemini API as fallback.

        Args:
            query: Original query string

        Returns:
            List of generated queries
        """
        try:
            # Get API key
            api_key = await self.gemini_key_manager.get_next_key()

            # Configure Gemini
            genai.configure(api_key=api_key)

            # Initialize model
            model = genai.GenerativeModel(
                model_name=Config.GEMINI_MODEL,
                generation_config=genai.types.GenerationConfig(
                    temperature=Config.GEMINI_TEMPERATURE,
                    max_output_tokens=Config.GEMINI_MAX_TOKENS,
                    top_p=Config.GEMINI_TOP_P,
                ),
                safety_settings={
                    HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
                    HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
                    HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
                    HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
                },
            )

            # Format prompt
            prompt = Config.get_prompt("rephrase_prompt", question=query)

            # Make request to Gemini
            response = await asyncio.to_thread(model.generate_content, prompt)
            response_text = response.text.strip()

            # Parse response into individual queries
            queries = self._parse_query_response(response_text)

            logger.debug(f"Gemini generated {len(queries)} queries")
            return queries

        except Exception as e:
            logger.error(f"Gemini query generation failed: {e}")
            raise

    def _parse_query_response(self, response_text: str) -> List[str]:
        """
        Parse LLM response into individual queries.

        Args:
            response_text: Raw response from LLM

        Returns:
            List of parsed queries
        """
        try:
            # Split by newlines and clean up
            lines = [line.strip() for line in response_text.split("\n") if line.strip()]

            # Remove any numbering, bullet points, or formatting
            queries = []
            for line in lines:
                # Remove common prefixes
                cleaned_line = re.sub(r"^[\d\.\-\*\+\s]*", "", line).strip()
                # Remove quotes if present
                cleaned_line = cleaned_line.strip("\"'")

                if cleaned_line and len(cleaned_line) > 10:  # Minimum query length
                    queries.append(cleaned_line)

            # Return up to 3 queries
            return queries[:3] if len(queries) >= 3 else queries

        except Exception as e:
            logger.error(f"Failed to parse query response: {e}")
            return []
